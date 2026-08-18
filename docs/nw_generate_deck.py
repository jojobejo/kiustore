import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Exact color constants matching user screenshots
C_NAVY_DARK = RGBColor(15, 23, 42)       # #0F172A
C_NAVY_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC
C_ORANGE_HEADER = RGBColor(234, 88, 12)  # #EA580C
C_BORDER = RGBColor(203, 213, 225)       # #CBD5E1
C_GREEN_STATUS = RGBColor(5, 150, 105)   # #059669
C_ORANGE_STATUS = RGBColor(217, 119, 6)  # #D97706
C_RED_STATUS = RGBColor(220, 38, 38)     # #DC2626
C_WHITE = RGBColor(255, 255, 255)
C_DARK_TEXT = RGBColor(15, 23, 42)
C_MUTED_TEXT = RGBColor(71, 85, 105)
C_CYAN_ACCENT = RGBColor(2, 132, 199)

def add_clean_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_WHITE
    bg.line.color.rgb = C_WHITE

# ==================== SLIDE 1: COVER ====================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = C_NAVY_DARK
bg1.line.color.rgb = C_NAVY_DARK

tb1 = s1.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.0))
tf1 = tb1.text_frame
tf1.word_wrap = True

p_badge = tf1.paragraphs[0]
p_badge.alignment = PP_ALIGN.CENTER
p_badge.text = "EXECUTIVE BALANCED SCORECARD REPORT 2026"
p_badge.font.size = Pt(11)
p_badge.font.bold = True
p_badge.font.color.rgb = C_ORANGE_HEADER

p_title = tf1.add_paragraph()
p_title.alignment = PP_ALIGN.CENTER
p_title.text = "KARISMA ONLINE // BALANCED SCORECARD"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = C_WHITE
p_title.space_before = Pt(8)

p_sub = tf1.add_paragraph()
p_sub.alignment = PP_ALIGN.CENTER
p_sub.text = "Evaluasi Kinerja Terpadu 9 Sisi Strategis Berbasis Sumber Data Tunggal (Single Source of Truth)"
p_sub.font.size = Pt(15)
p_sub.font.color.rgb = RGBColor(148, 163, 184)
p_sub.space_before = Pt(8)

p_meta = tf1.add_paragraph()
p_meta.alignment = PP_ALIGN.CENTER
p_meta.text = "PT. KARISMA INDOAGRO UNIVERSAL  •  PERIODE 2026"
p_meta.font.size = Pt(11)
p_meta.font.bold = True
p_meta.font.color.rgb = C_CYAN_ACCENT
p_meta.space_before = Pt(20)

kpi_data = [
    ("94.2%", "Rata-Rata 9 BSC"),
    ("< 5 Detik", "Otomasi BRIVA"),
    ("100%", "Kepatuhan App Store"),
    ("Kecamatan", "Akurasi Logistik")
]
for i, (val, lbl) in enumerate(kpi_data):
    k_left = Inches(1.5 + i * 2.65)
    k_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, k_left, Inches(5.2), Inches(2.4), Inches(1.3))
    k_card.fill.solid()
    k_card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    k_card.line.color.rgb = RGBColor(51, 65, 85)
    
    k_tb = s1.shapes.add_textbox(k_left, Inches(5.3), Inches(2.4), Inches(1.1))
    k_tf = k_tb.text_frame
    kp0 = k_tf.paragraphs[0]
    kp0.alignment = PP_ALIGN.CENTER
    kp0.text = val
    kp0.font.size = Pt(18)
    kp0.font.bold = True
    kp0.font.color.rgb = RGBColor(56, 189, 248)
    
    kp1 = k_tf.add_paragraph()
    kp1.alignment = PP_ALIGN.CENTER
    kp1.text = lbl
    kp1.font.size = Pt(9.5)
    kp1.font.bold = True
    kp1.font.color.rgb = RGBColor(148, 163, 184)

# ==================== SLIDE 2: MASTER DASHBOARD 9 BSC ====================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_bg(s2)

# Top Bar Header
tb_top = s2.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.133), Inches(0.4))
tf_top = tb_top.text_frame
p_tl = tf_top.paragraphs[0]
p_tl.text = "KARISMA ONLINE // BALANCED SCORECARD (BSC) AUDIT REPORT              SUMBER DATA TERVERIFIKASI"
p_tl.font.size = Pt(9.5)
p_tl.font.bold = True
p_tl.font.color.rgb = C_DARK_TEXT

# Main Title
tb_m = s2.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.133), Inches(0.7))
tf_m = tb_m.text_frame
pm0 = tf_m.paragraphs[0]
pm0.text = "DASHBOARD EKSEKUTIF: 9 BALANCED SCORECARDS (BSC) KARISMA"
pm0.font.size = Pt(16)
pm0.font.bold = True
pm0.font.color.rgb = C_DARK_TEXT

pm1 = tf_m.add_paragraph()
pm1.text = "Evaluasi Kinerja Terpadu Berbasis Sumber Data Tunggal (Single Source of Truth)"
pm1.font.size = Pt(10)
pm1.font.color.rgb = C_MUTED_TEXT

# Master Table (10 rows, 6 cols)
t_shape2 = s2.shapes.add_table(10, 6, Inches(0.6), Inches(1.2), Inches(12.133), Inches(3.6))
t2 = t_shape2.table
t2.columns[0].width = Inches(0.5)
t2.columns[1].width = Inches(2.6)
t2.columns[2].width = Inches(2.8)
t2.columns[3].width = Inches(3.833)
t2.columns[4].width = Inches(1.0)
t2.columns[5].width = Inches(1.4)

headers2 = ["No", "Perspektif BSC", "Unit Penanggung Jawab (Owner)", "Measure Lead Utama", "Overall", "Indikator Kinerja"]
for j, h in enumerate(headers2):
    cell = t2.cell(0, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_ORANGE_HEADER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    if j == 0 or j >= 4:
        p.alignment = PP_ALIGN.CENTER

master_rows = [
    ("1", "BSC Karisma Online", "Tim Digital & IT Karisma Online", "Status Kesiapan API, Test Coverage, Kepatuhan App Store", "95%", "■ On Progress"),
    ("2", "BSC Sisi Customer", "Tim Customer Experience & Komersial", "Indeks Rating Layanan Sales, Retensi Mitra, Siklus Order", "94%", "■ On Progress"),
    ("3", "BSC Divisi Sales", "Kepala Divisi Penjualan & Distribusi", "Utilisasi Limit Kredit Toko, Rasio Toko Aktif Order", "89%", "■ Perlu Intervensi"),
    ("4", "BSC Logistik & Warehouse", "Tim Logistik & Operasional Gudang", "SLA Pengemasan (<2 Jam), Akurasi Ongkir Kecamatan", "92%", "■ On Progress"),
    ("5", "BSC Sisi Multiplatform", "Tim Mobile Engineering & UI/UX", "App Store Review Clearance, Android Crash Rate (<0.1%)", "96%", "■ On Progress"),
    ("6", "BSC Sisi Perusahaan", "Dewan Direksi & Tim Manajemen", "Days Sales Outstanding (DSO), Efisiensi Biaya per Order", "94%", "■ On Progress"),
    ("7", "BSC Payment & Keuangan", "Divisi Keuangan & Treasury", "Adopsi BRIVA, Kecepatan Rekonsiliasi Kas, Rasio Over-Limit", "96%", "■ On Progress"),
    ("8", "BSC Sisi Teknis", "Tim Lead Backend Engineering", "API Latency (<200ms), Error Rate (<0.01%), ACID Transaksi", "95%", "■ On Progress"),
    ("9", "BSC Sisi Keamanan", "Tim Security & Compliance", "Zero Critical Vulnerability, Token Expiry, Audit Non-PII", "97%", "■ On Progress")
]

for i, row in enumerate(master_rows):
    row_idx = i + 1
    for j, val in enumerate(row):
        cell = t2.cell(row_idx, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY_LIGHT if i % 2 == 1 else C_WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(8.5)
        p.font.color.rgb = C_DARK_TEXT
        if j == 0:
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
        elif j == 1:
            p.font.bold = True
        elif j == 4:
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            cell.fill.solid()
            if "89%" in val:
                cell.fill.fore_color.rgb = C_RED_STATUS
            else:
                cell.fill.fore_color.rgb = C_ORANGE_STATUS
            p.font.color.rgb = C_WHITE
        elif j == 5:
            p.font.size = Pt(8)
            p.font.bold = True
            if "Intervensi" in val:
                p.font.color.rgb = C_RED_STATUS
            else:
                p.font.color.rgb = C_ORANGE_STATUS

# 4-Grid Integrity Box below
tb_int_t = s2.shapes.add_textbox(Inches(0.6), Inches(4.95), Inches(12.133), Inches(0.3))
p_it = tb_int_t.text_frame.paragraphs[0]
p_it.text = "PRINSIP INTEGRITAS DATA & SUMBER DOKUMENTASI RESMI"
p_it.font.size = Pt(10)
p_it.font.bold = True
p_it.font.color.rgb = C_DARK_TEXT

grid_data = [
    ("1. Arsitektur & Kepatuhan Platform (95%)", "Terverifikasi pada controller mobile CodeIgniter (Mobile.php) dan arsitektur iOS Swift/Android Kotlin. Lolos Apple Review 5.1.1 melalui audit tabel mobile_account_deletions."),
    ("2. Finansial & Limit Kredit Pelanggan (94%)", "Data tagihan berjalan dan limit kredit dibaca langsung dari customers.max_credit dan agregasi tabel orders serta view database v_products."),
    ("3. Kesiapan Divisi Sales & Distribusi (89%)", "Pemetaan akun toko binaan aktif di backend (customers.salesman_id). Tantangan utama terletak pada legalitas skema komisi teritori digital untuk mendorong adopsi lapangan."),
    ("4. Ekosistem Pembayaran & Virtual Account (96%)", "Kanal BRIVA aktif via library Brivaws pada tabel briva_api dengan auto-update status pembayaran instan ke pesanan lunas (status = 10).")
]

for idx, (g_title, g_desc) in enumerate(grid_data):
    gx = Inches(0.6 + (idx % 2) * 6.15)
    gy = Inches(5.3 + (idx // 2) * 0.95)
    
    g_card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gx, gy, Inches(5.98), Inches(0.85))
    g_card.fill.solid()
    g_card.fill.fore_color.rgb = C_NAVY_LIGHT
    g_card.line.color.rgb = C_BORDER
    
    g_tb = s2.shapes.add_textbox(gx + Inches(0.1), gy + Inches(0.05), Inches(5.78), Inches(0.75))
    g_tf = g_tb.text_frame
    g_tf.word_wrap = True
    
    gp0 = g_tf.paragraphs[0]
    gp0.text = g_title
    gp0.font.size = Pt(9)
    gp0.font.bold = True
    gp0.font.color.rgb = C_DARK_TEXT
    
    gp1 = g_tf.add_paragraph()
    gp1.text = g_desc
    gp1.font.size = Pt(8)
    gp1.font.color.rgb = C_MUTED_TEXT

# ==================== HELPER FUNCTION FOR 9 BSC DETAIL SLIDES WITH MERGED MUST WIN ====================
def create_exact_bsc_slide(slide_num, bsc_title, obj_desc, meta_dict, initiatives, total_score, p_verif, t_isu, r_mitigasi):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_clean_bg(s)
    
    # 1. Header Banner
    h_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.3), Inches(12.133), Inches(0.45))
    h_box.fill.solid()
    h_box.fill.fore_color.rgb = C_NAVY_DARK
    h_box.line.color.rgb = C_NAVY_DARK
    h_tb = s.shapes.add_textbox(Inches(0.7), Inches(0.32), Inches(11.933), Inches(0.4))
    hp = h_tb.text_frame.paragraphs[0]
    hp.alignment = PP_ALIGN.CENTER
    hp.text = f"BALANCED SCORECARD: {bsc_title.upper()}"
    hp.font.size = Pt(11.5)
    hp.font.bold = True
    hp.font.color.rgb = C_WHITE
    
    # 2. Objective Row
    o_lbl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.6))
    o_lbl.fill.solid()
    o_lbl.fill.fore_color.rgb = C_NAVY_DARK
    o_lbl.line.color.rgb = C_BORDER
    o_ltb = s.shapes.add_textbox(Inches(0.65), Inches(0.85), Inches(2.1), Inches(0.4))
    olp = o_ltb.text_frame.paragraphs[0]
    olp.alignment = PP_ALIGN.CENTER
    olp.text = "Objective Description"
    olp.font.size = Pt(9.5)
    olp.font.bold = True
    olp.font.color.rgb = C_WHITE
    
    o_val = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.8), Inches(0.75), Inches(9.933), Inches(0.6))
    o_val.fill.solid()
    o_val.fill.fore_color.rgb = C_WHITE
    o_val.line.color.rgb = C_BORDER
    o_vtb = s.shapes.add_textbox(Inches(2.9), Inches(0.78), Inches(9.733), Inches(0.54))
    o_vtf = o_vtb.text_frame
    o_vtf.word_wrap = True
    ovp = o_vtf.paragraphs[0]
    ovp.text = obj_desc
    ovp.font.size = Pt(8.5)
    ovp.font.color.rgb = C_DARK_TEXT
    
    # 3. Meta Bar
    m_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(12.133), Inches(0.4))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = C_NAVY_LIGHT
    m_box.line.color.rgb = C_BORDER
    m_tb = s.shapes.add_textbox(Inches(0.7), Inches(1.38), Inches(11.933), Inches(0.35))
    mp = m_tb.text_frame.paragraphs[0]
    mp.text = f"Owner: {meta_dict['owner']}    |    Measure lead: {meta_dict['lead']}    |    Frequency: {meta_dict['freq']}    |    Overall Status: {meta_dict['status']}    |    {meta_dict['date']}"
    mp.font.size = Pt(9)
    mp.font.bold = True
    mp.font.color.rgb = C_DARK_TEXT
    
    # 4. Table (rows = len(initiatives) + 2 for header and total)
    rows = len(initiatives) + 2
    t_shape = s.shapes.add_table(rows, 4, Inches(0.6), Inches(1.75), Inches(12.133), Inches(3.2))
    t = t_shape.table
    t.columns[0].width = Inches(2.8)
    t.columns[1].width = Inches(4.833)
    t.columns[2].width = Inches(3.3)
    t.columns[3].width = Inches(1.2)
    
    headers = ["Must Win", "Key initiatives", "Dasar Verifikasi / Sumber Data", "Status"]
    for j, h in enumerate(headers):
        cell = t.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_ORANGE_HEADER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        if j == 3:
            p.alignment = PP_ALIGN.CENTER
            
    # Grouping logic for Must Win
    groups = []
    current_mw = None
    current_start = 0

    for i, (mw, init, audit, st) in enumerate(initiatives):
        if current_mw is None:
            current_mw = mw
            current_start = i
        elif mw != current_mw:
            groups.append((current_start, i - 1, current_mw))
            current_mw = mw
            current_start = i
    if current_mw is not None:
        groups.append((current_start, len(initiatives) - 1, current_mw))
            
    for i, (mw, init, audit, st) in enumerate(initiatives):
        row_idx = i + 1
        cell0 = t.cell(row_idx, 0)
        cell0.fill.solid()
        cell0.fill.fore_color.rgb = C_WHITE
        cell0.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = cell0.text_frame.paragraphs[0]
        p0.text = mw
        p0.font.size = Pt(8.5)
        p0.font.bold = True
        p0.font.color.rgb = C_DARK_TEXT
        
        cell1 = t.cell(row_idx, 1)
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = C_WHITE
        cell1.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = cell1.text_frame.paragraphs[0]
        p1.text = init
        p1.font.size = Pt(8.5)
        p1.font.color.rgb = C_DARK_TEXT
        
        cell2 = t.cell(row_idx, 2)
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = C_WHITE
        cell2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = cell2.text_frame.paragraphs[0]
        p2.text = audit
        p2.font.size = Pt(8)
        p2.font.italic = True
        p2.font.color.rgb = C_MUTED_TEXT
        
        cell3 = t.cell(row_idx, 3)
        cell3.fill.solid()
        if "100%" in st:
            cell3.fill.fore_color.rgb = C_GREEN_STATUS
        elif "95%" in st or "90%" in st:
            cell3.fill.fore_color.rgb = C_ORANGE_STATUS
        else:
            cell3.fill.fore_color.rgb = C_RED_STATUS
        cell3.vertical_anchor = MSO_ANCHOR.MIDDLE
        p3 = cell3.text_frame.paragraphs[0]
        p3.text = st
        p3.alignment = PP_ALIGN.CENTER
        p3.font.size = Pt(9)
        p3.font.bold = True
        p3.font.color.rgb = C_WHITE
        
    # MERGE CELLS IN PPTX FOR MUST WIN
    for g_start, g_end, mw_text in groups:
        if g_end > g_start:
            c_top = t.cell(g_start + 1, 0)
            c_bot = t.cell(g_end + 1, 0)
            c_top.merge(c_bot)
            c_top.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c_top.text_frame.paragraphs[0]
            p.text = mw_text
            p.font.size = Pt(8.5)
            p.font.bold = True
            p.font.color.rgb = C_DARK_TEXT
        
    # Total Row
    tot_idx = rows - 1
    cell_tot_lbl = t.cell(tot_idx, 0)
    cell_tot_lbl.fill.solid()
    cell_tot_lbl.fill.fore_color.rgb = C_NAVY_DARK
    cell_tot_lbl.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    t.cell(tot_idx, 1).fill.solid()
    t.cell(tot_idx, 1).fill.fore_color.rgb = C_NAVY_DARK
    t.cell(tot_idx, 2).fill.solid()
    t.cell(tot_idx, 2).fill.fore_color.rgb = C_NAVY_DARK
    
    p_tot_lbl = t.cell(tot_idx, 2).text_frame.paragraphs[0]
    p_tot_lbl.alignment = PP_ALIGN.RIGHT
    p_tot_lbl.text = "Total Capaian Kinerja Terverifikasi (Overall Score)"
    p_tot_lbl.font.size = Pt(9)
    p_tot_lbl.font.bold = True
    p_tot_lbl.font.color.rgb = C_WHITE
    
    cell_tot_v = t.cell(tot_idx, 3)
    cell_tot_v.fill.solid()
    if "89%" in total_score:
        cell_tot_v.fill.fore_color.rgb = C_RED_STATUS
    else:
        cell_tot_v.fill.fore_color.rgb = C_ORANGE_HEADER
    cell_tot_v.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_tot_v = cell_tot_v.text_frame.paragraphs[0]
    p_tot_v.alignment = PP_ALIGN.CENTER
    p_tot_v.text = total_score
    p_tot_v.font.size = Pt(10)
    p_tot_v.font.bold = True
    p_tot_v.font.color.rgb = C_WHITE

    # 5. Legend Strip
    leg_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.95), Inches(12.133), Inches(0.3))
    leg_box.fill.solid()
    leg_box.fill.fore_color.rgb = C_NAVY_DARK
    leg_box.line.color.rgb = C_BORDER
    leg_tb = s.shapes.add_textbox(Inches(0.7), Inches(4.97), Inches(11.933), Inches(0.25))
    lp = leg_tb.text_frame.paragraphs[0]
    lp.text = "■ >= 100% (Tercapai / Teruji Lulus Penuh)      ■ On Progress <100% - 90% (Live Staging / Dalam Pengawalan)      ■ Belum Start / <90% (Perlu Kebijakan Manajemen)"
    lp.font.size = Pt(8.5)
    lp.font.bold = True
    lp.font.color.rgb = C_WHITE
    
    # 6. Grounded 3-Column Box below
    g_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.3), Inches(12.133), Inches(1.85))
    g_box.fill.solid()
    g_box.fill.fore_color.rgb = C_WHITE
    g_box.line.color.rgb = C_BORDER
    
    cols_data = [
        ("PENCAPAIAN TERVERIFIKASI", C_GREEN_STATUS, p_verif),
        ("TANTANGAN & ISU LAPANGAN", C_ORANGE_HEADER, t_isu),
        ("RENCANA AKSI & MITIGASI RISIKO", C_CYAN_ACCENT, r_mitigasi)
    ]
    for c_idx, (col_title, col_color, col_items) in enumerate(cols_data):
        cx = Inches(0.7 + c_idx * 4.05)
        c_tb = s.shapes.add_textbox(cx, Inches(5.35), Inches(3.9), Inches(1.75))
        c_tf = c_tb.text_frame
        c_tf.word_wrap = True
        
        cp0 = c_tf.paragraphs[0]
        cp0.text = col_title
        cp0.font.size = Pt(9)
        cp0.font.bold = True
        cp0.font.color.rgb = col_color
        
        for item in col_items:
            cp = c_tf.add_paragraph()
            cp.text = f"• {item}"
            cp.font.size = Pt(7.5)
            cp.font.color.rgb = C_DARK_TEXT
            cp.space_before = Pt(2)
            
    return s

# ==================== SLIDE 3: BSC 1 - APPS / KARISMA ONLINE ====================
create_exact_bsc_slide(
    3,
    "KARISMA ONLINE (PLATFORM & DIGITAL ECOSYSTEM)",
    "Digitalisasi menyeluruh kanal pemesanan B2B/B2C, integrasi monitoring limit kredit & penagihan virtual account (BRIVA), percepatan pemenuhan order, serta kesiapan distribusi multi-platform yang aman dan patuh regulasi.",
    {"owner": "Tim Digital & IT Karisma Online", "lead": "Status Kesiapan API, Test Coverage, Kepatuhan App Store", "freq": "Bulanan", "status": "95%", "date": "31-Agu-26"},
    [
        ("1. KEANDALAN AKSES & PENGALAMAN PENGGUNA", "1.1 Akses pemesanan mandiri 24/7 (Aplikasi iOS Swift, Android & Web KiuStore)", "Terverifikasi: MainTabView, Catalog, Cart, Checkout", "100%"),
        ("1. KEANDALAN AKSES & PENGALAMAN PENGGUNA", "1.2 Mode jelajah katalog produk & promo tanpa kewajiban login awal (Guest Browsing)", "Terverifikasi: GuestAccessTests & LoginView", "100%"),
        ("1. KEANDALAN AKSES & PENGALAMAN PENGGUNA", "1.3 Sinkronisasi harga grosir & stok produk real-time via API /api/v1", "Terverifikasi: APIEndpoint & Staging API", "95%"),
        ("2. INTEGRASI FINANSIAL & TRANSAKSI", "2.1 Integrasi pembayaran Virtual Account instan BRIVA via Brivaws", "Terverifikasi: MobileBrivaPaymentAPI & OrderTests", "95%"),
        ("2. INTEGRASI FINANSIAL & TRANSAKSI", "2.2 Visibilitas limit kredit (max_credit) & tagihan berjalan real-time bagi mitra", "Terverifikasi: CustomerFinanceTests & GET /customer/finance-summary", "100%"),
        ("2. INTEGRASI FINANSIAL & TRANSAKSI", "2.3 Otomasi validasi persetujuan pesanan sesuai plafon kredit toko (payment_method = 1)", "Terverifikasi: CheckoutViewModelTests", "100%"),
        ("3. KECEPATAN LOGISTIK & ORDER", "3.1 Pemrosesan alur pesanan dari checkout hingga pengiriman (Order SLA)", "Terverifikasi: Alur status pesanan 1, 2, 3, 4, 5", "90%"),
        ("3. KECEPATAN LOGISTIK & ORDER", "3.2 Transparansi status kirim, ongkir (mobile_shipping_quotes) & faktur digital", "Terverifikasi: ShippingTests & InvoiceView", "95%"),
        ("3. KECEPATAN LOGISTIK & ORDER", "3.3 Fitur pemesanan cepat rutin 1-sentuhan (Quick Re-Order)", "Status: Dalam pengembangan modul", "85%"),
        ("4. KEPATUHAN REGULASI & DATA", "4.1 Kepatuhan 100% standar Apple App Review & Privasi (Hapus Akun DELETE /account)", "Terverifikasi: AppReviewResolution & ProductionSafetyTests", "100%"),
        ("4. KEPATUHAN REGULASI & DATA", "4.2 Layanan bantuan pelanggan interaktif terintegrasi (In-App Chat Support)", "Terverifikasi: ChatTests & ChatView", "95%"),
        ("4. KEPATUHAN REGULASI & DATA", "4.3 Pembentukan duta digital (Digital Champion) di cabang untuk pendampingan toko", "Status: Menunggu standardisasi SOP cabang", "80%")
    ],
    "95%",
    ["Lolos 100% kepatuhan Apple Guideline 5.1.1 (Penghapusan akun terenkripsi SHA-256 pada mobile_account_deletions).", "Integrasi limit kredit real-time (customers.max_credit) aktif mencegah pesanan over-limit."],
    ["Penyesuaian pemahaman alur digital bagi toko-toko konvensional di area perintis.", "Alur retur barang rusak masih mengandalkan verifikasi fisik manual di depo/gudang."],
    ["Meluncurkan fitur Quick Re-Order 1-klik untuk memudahkan pesanan paket mingguan toko mitra.", "Menetapkan target penyiapan barang gudang seragam (SLA 2 jam) dan modul unggah foto retur digital."]
)

# ==================== SLIDE 4: BSC 2 - SISI CUSTOMER ====================
create_exact_bsc_slide(
    4,
    "SISI CUSTOMER (CUSTOMER EXPERIENCE & MITRA ACCESS)",
    "Memberikan kemudahan pemesanan mandiri 24/7, fleksibilitas tingkatan harga sesuai level mitra, transparansi kupon diskon, layanan interaktif responsif, dan keterbukaan ulasan kepuasan produk.",
    {"owner": "Tim Customer Experience & Komersial", "lead": "Indeks Rating Layanan Sales, Retensi Mitra, Siklus Order", "freq": "Bulanan", "status": "94%", "date": "31-Agu-26"},
    [
        ("1. KEMUDAHAN AKSES & BELANJA", "1.1 Mode Guest Browsing katalog produk tanpa kewajiban login awal", "Terverifikasi: Endpoint GET /api/v1/products & UI Katalog", "100%"),
        ("1. KEMUDAHAN AKSES & BELANJA", "1.2 Manajemen data profil toko, kontak, dan multi-alamat pengiriman", "Terverifikasi: Endpoint PUT /api/v1/profile & Tabel customers", "100%"),
        ("1. KEMUDAHAN AKSES & BELANJA", "1.3 Sistem kupon promo potongan harga transaksi belanja", "Terverifikasi: Tabel coupons & validasi Shop.php", "95%"),
        ("2. LAYANAN INTERAKTIF & CHAT", "2.1 Saluran pesan langsung Customer Service 2 arah (In-App Live Chat)", "Terverifikasi: Endpoint /api/v1/messages & Tabel message", "95%"),
        ("2. LAYANAN INTERAKTIF & CHAT", "2.2 Sistem ulasan dan rating bintang kualitas produk pasca pesanan selesai", "Terverifikasi: Endpoint POST /api/v1/orders/{id}/complete", "95%"),
        ("2. LAYANAN INTERAKTIF & CHAT", "2.3 Notifikasi pengingat pembayaran & status pesanan instan", "Status: Integrasi gateway notifikasi pesan", "85%"),
        ("3. TRANSPARANSI & LOYALITAS", "3.1 Akses riwayat faktur, rincian biaya, dan status pengiriman real-time", "Terverifikasi: Endpoint GET /api/v1/orders & Tabel orders", "100%"),
        ("3. TRANSPARANSI & LOYALITAS", "3.2 Program poin loyalitas dan reward kuota diskon toko aktif", "Status: Tahap perancangan skema reward Q4 2026", "80%")
    ],
    "94%",
    ["Katalog terbuka tanpa login memudahkan eksplorasi produk.", "Live Chat CS aktif menghubungkan kios langsung ke operator kantor pusat."],
    ["Kios pedesaan masih terbiasa menelepon pribadi salesman."],
    ["Menyediakan video panduan ringkas 1 menit dan voucher promo belanja perdana via aplikasi."]
)

# ==================== SLIDE 5: BSC 3 - SISI SALES ====================
create_exact_bsc_slide(
    5,
    "SISI SALES (SALES FORCE ENABLEMENT & DISTRIBUSI)",
    "Memberdayakan tenaga penjual lapangan, mengoptimalkan pembagian wilayah kios binaan, mengawal kepatuhan limit piutang, dan mempercepat rekonsiliasi order digital.",
    {"owner": "Kepala Divisi Penjualan & Distribusi", "lead": "Utilisasi Limit Kredit, Rasio Toko Aktif, Rating Layanan Sales", "freq": "Bulanan", "status": "89%", "date": "31-Agu-26"},
    [
        ("1. PEMETAAN WILAYAH & TOKO", "1.1 Penugasan kios binaan per salesman lapangan di basis data terpusat", "Terverifikasi: Field customers.salesman_id & Admin Salesman", "100%"),
        ("1. PEMETAAN WILAYAH & TOKO", "1.2 Visibilitas riwayat transaksi dan status pembayaran toko binaan", "Terverifikasi: Modul Admin Salesman API & Controller Orders", "95%"),
        ("1. PEMETAAN WILAYAH & TOKO", "1.3 Perencanaan rute kunjungan sales terintegrasi riwayat order kios", "Status: Dalam perancangan modul territory route", "80%"),
        ("2. PENGAWALAN PLAFON KREDIT", "2.1 Visibilitas limit kredit (max_credit) toko binaan saat pemesanan", "Terverifikasi: Model Mobile_api_model & Admin Piutang", "95%"),
        ("2. PENGAWALAN PLAFON KREDIT", "2.2 Alur validasi persetujuan pesanan tempo digital oleh supervisor sales", "Terverifikasi: Status alur piutang di Piutang.php", "90%"),
        ("2. PENGAWALAN PLAFON KREDIT", "2.3 Skema insentif komisi teritori digital untuk mendorong adopsi mobile", "Status: Menunggu legalitas formal skema komisi direksi", "80%"),
        ("3. PENDAMPINGAN & KUALITAS", "3.1 Program pendampingan instalasi aplikasi ke kios mitra perintis", "Terverifikasi: Pelaksanaan onboarding tim cabang", "85%"),
        ("3. PENDAMPINGAN & KUALITAS", "3.2 Sistem penilaian kinerja pelayanan salesman dari ulasan toko", "Terverifikasi: Modul Rating.php & Tabel reviews", "95%")
    ],
    "89%",
    ["Struktur database mengunci relasi toko ke salesman_id.", "Limit kredit toko terkontrol otomatis mencegah resiko over-limit."],
    ["Kekhawatiran tim sales bahwa order mandiri aplikasi akan memotong komisi penjualan mereka."],
    ["Mengesahkan SK Direksi: Komisi pesanan digital tetap 100% dialokasikan ke salesman pembina."]
)

# ==================== SLIDE 6: BSC 4 - SISI LOGISTIK & WAREHOUSE ====================
create_exact_bsc_slide(
    6,
    "SISI LOGISTIK & WAREHOUSE (SUPPLY CHAIN & SLA)",
    "Menjamin presisi perhitungan ongkos kirim hingga tingkat kecamatan, efisiensi alur pengemasan gudang (Order SLA), mitigasi produk lambat bergerak (deadstock), dan penyediaan armada muatan besar.",
    {"owner": "Tim Logistik & Operasional Gudang", "lead": "SLA Pengemasan (<2 Jam), Akurasi Ongkir, Rasio Deadstock", "freq": "Harian", "status": "92%", "date": "31-Agu-26"},
    [
        ("1. PRESISI LOGISTIK WILAYAH", "1.1 Integrasi pembacaan tarif ekspedisi tingkat kecamatan (RajaOngkir Pro API)", "Terverifikasi: Endpoint shipping/destination & Rajaongkir.php", "100%"),
        ("1. PRESISI LOGISTIK WILAYAH", "1.2 Penguncian kuotasi tarif ongkos kirim 30 menit saat proses checkout", "Terverifikasi: Tabel mobile_shipping_quotes & expiry lock", "100%"),
        ("1. PRESISI LOGISTIK WILAYAH", "1.3 Kalkulasi otomatis akumulasi berat produk satuan botol/pcs dan dus/karton", "Terverifikasi: Model validasi berat & product_unit_value", "100%"),
        ("2. KECEPATAN PEMENUHAN GUDANG", "2.1 Antrean status alur pesanan terstruktur (Verifikasi -> Kemas -> Kirim)", "Terverifikasi: Controller Pengiriman.php & Admin Orders", "95%"),
        ("2. KECEPATAN PEMENUHAN GUDANG", "2.2 Standarisasi target penyiapan barang gudang seragam (SLA < 2 Jam)", "Status: Penyelarasan SOP shift tim gudang cabang", "85%"),
        ("2. KECEPATAN PEMENUHAN GUDANG", "2.3 Skema pengiriman armada truk internal Karisma untuk muatan tonase besar", "Status: Penentuan tarif flat rute rutin armada internal", "80%"),
        ("3. PENGENDALIAN DEADSTOCK", "3.1 Deteksi dini produk slow-moving (> 1 tahun pergerakan < 5% & retur)", "Terverifikasi: Catatan teknis CATATAN PENTING.txt #DEADSTOCK", "90%"),
        ("3. PENGENDALIAN DEADSTOCK", "3.2 Otomasi alur diskon kilat (flash sale) cuci gudang produk deadstock", "Status: Tahap sinkronisasi formula promo staging", "80%")
    ],
    "92%",
    ["Penguncian ongkir 30 menit (mobile_shipping_quotes) berhasil mengeliminasi selisih tarif.", "Alur status pesanan gudang terintegrasi langsung dengan nomor resi."],
    ["Tarif ekspedisi kurir reguler tidak ekonomis untuk pesanan pupuk tonase besar (drum)."],
    ["Mengaktifkan opsi armada truk internal Karisma dengan tarif flat terjangkau rute mingguan."]
)

# ==================== SLIDE 7: BSC 5 - SISI MULTIPLATFORM ====================
create_exact_bsc_slide(
    7,
    "SISI MULTIPLATFORM (TRI-PLATFORM IOS, ANDROID, WEB)",
    "Menghadirkan performa aplikasi native berkecepatan tinggi, ukuran aplikasi hemat kuota, antarmuka responsif di ponsel/tablet, dan stabilitas back-office web portal.",
    {"owner": "Tim Mobile Engineering & UI/UX", "lead": "App Store Review Clearance, Android Crash Rate (<0.1%)", "freq": "Bulanan", "status": "96%", "date": "31-Agu-26"},
    [
        ("1. EKOSISTEM IOS (APP STORE)", "1.1 Arsitektur Swift Native (SwiftUI & UIKit) berkinerja tinggi", "Terverifikasi: Source build iOS Karisma Online", "100%"),
        ("1. EKOSISTEM IOS (APP STORE)", "1.2 Kepatuhan Apple Review Guideline 5.1.1(v) (Penghapusan Akun Mandiri)", "Terverifikasi: Dokumen docs/APP_REVIEW_RESOLUTION_20260728.md", "100%"),
        ("1. EKOSISTEM IOS (APP STORE)", "1.3 Tata letak antarmuka adaptif layar iPad dengan tombol navigasi jelas", "Terverifikasi: Kepatuhan review UI iPad resolution", "95%"),
        ("2. EKOSISTEM ANDROID (PLAY STORE)", "2.1 Arsitektur Kotlin/Java Native (MVVM Pattern) responsif", "Terverifikasi: Source code kiustore_apps module", "100%"),
        ("2. EKOSISTEM ANDROID (PLAY STORE)", "2.2 Optimasi ukuran installer APK hemat kuota internet (< 15MB)", "Terverifikasi: Build packaging & resource shrinking", "95%"),
        ("2. EKOSISTEM ANDROID (PLAY STORE)", "2.3 Kompresi otomatis foto bukti transfer bank sebelum diunggah", "Terverifikasi: Controller Mobile.php payment_picture_name", "95%"),
        ("3. PORTAL BACK-OFFICE WEB", "3.1 Dasbor operasional web AdminLTE dengan manajemen modul lengkap", "Terverifikasi: Module application/modules/admin", "100%"),
        ("3. PORTAL BACK-OFFICE WEB", "3.2 Sinkronisasi status transaksi lintas platform (iOS, Android, Web)", "Terverifikasi: Database MariaDB InnoDB ACID transactions", "95%")
    ],
    "96%",
    ["Lulus 100% audit kepatuhan Apple Review (DELETE /account non-PII & guest browsing).", "APK Android berukuran ringan di bawah 15MB hemat kuota bagi mitra pedesaan."],
    ["Variasi performa smartphone Android tipe lama milik sebagian kios mitra binaan."],
    ["Mengoptimalkan caching lokal dan kompresi gambar otomatis pada sisi perangkat."]
)

# ==================== SLIDE 8: BSC 6 - SISI PERUSAHAAN ====================
create_exact_bsc_slide(
    8,
    "SISI PERUSAHAAN (GOVERNANCE & CORPORATE FINANCE)",
    "Mengakselerasi perputaran modal kerja, menekan Days Sales Outstanding (DSO), menjamin kepatuhan audit perpajakan, dan meningkatkan efisiensi biaya operasional per order.",
    {"owner": "Dewan Direksi & Tim Manajemen Eksekutif", "lead": "Days Sales Outstanding (DSO), Efisiensi Biaya, Kepatuhan Audit", "freq": "Bulanan", "status": "94%", "date": "31-Agu-26"},
    [
        ("1. ARUS KAS & MODAL KERJA", "1.1 Pemotongan waktu settlement pembayaran transfer dari 2-4 jam ke < 5 detik", "Terverifikasi: Otomasi webhook callback BRIVA", "100%"),
        ("1. ARUS KAS & MODAL KERJA", "1.2 Penurunan risiko piutang tak tertagih via validasi limit plafon kredit sistem", "Terverifikasi: Logika max_credit di checkout mobile", "95%"),
        ("1. ARUS KAS & MODAL KERJA", "1.3 Visibilitas dasbor eksekutif pergerakan omset harian dan bulanan", "Terverifikasi: Controller Dashboard.php & Report.php", "95%"),
        ("2. TATA KELOLA HUKUM & AUDIT", "2.1 Pemisahan data identitas pribadi (non-PII) tanpa menghapus faktur transaksi", "Terverifikasi: Model Mobile_api_model.php delete_account", "100%"),
        ("2. TATA KELOLA HUKUM & AUDIT", "2.2 Perlindungan integritas nomor faktur dan pesanan untuk audit perpajakan", "Terverifikasi: Relasi tabel orders, order_items, payments", "100%"),
        ("2. TATA KELOLA HUKUM & AUDIT", "2.3 Pencatatan audit trail lengkap seluruh aktivitas transaksi dan pembayaran", "Terverifikasi: Tabel briva_api, mobile_account_deletions", "95%"),
        ("3. SKALABILITAS KORPORASI", "3.1 Efisiensi biaya operasional per transaksi pemesanan", "Terverifikasi: Otomasi sistem mengurangi beban kasir", "90%"),
        ("3. SKALABILITAS KORPORASI", "3.2 Standardisasi SOP operasional cabang untuk ekspansi digital", "Status: Penyusunan panduan operasional cabang terpadu", "85%")
    ],
    "94%",
    ["Integritas faktur dan nomor pesanan terlindungi untuk keperluan audit pajak resmi.", "Arus kas bertumbuh lebih cepat berkat eliminasi jeda verifikasi manual bank."],
    ["Masa transisi staf administrasi cabang dari kebiasaan input data konvensional."],
    ["Menjadwalkan workshop pelatihan sistem digital terpadu bagi staf administrasi cabang."]
)

# ==================== SLIDE 9: BSC 7 - SISI KEUANGAN & PAYMENT ====================
create_exact_bsc_slide(
    9,
    "SISI KEUANGAN & PAYMENT (FINTECH & PRICING ENGINE)",
    "Otomasi penerimaan pembayaran Virtual Account (BRIVA), perlindungan skema harga bertingkat, kepastian rekonsiliasi kas harian, dan pencegahan transaksi over-limit.",
    {"owner": "Divisi Keuangan & Treasury Karisma", "lead": "Adopsi BRIVA, Kecepatan Rekonsiliasi Kas, Rasio Over-Limit", "freq": "Harian", "status": "96%", "date": "31-Agu-26"},
    [
        ("1. OTOMASI FINTECH BRIVA", "1.1 Integrasi library Brivaws standar SNAP API BRI (/transfer-va/create-va)", "Terverifikasi: Controller Brivawsapi.php & key/private.pem", "95%"),
        ("1. OTOMASI FINTECH BRIVA", "1.2 Penerbitan Dynamic VA (91118 + No HP) dengan batas kedaluwarsa 15 menit", "Terverifikasi: CATATAN PENTING.txt (V) & Tabel briva_api", "100%"),
        ("1. OTOMASI FINTECH BRIVA", "1.3 Webhook callback otomatis memperbarui status pesanan menjadi Lunas (status = 10)", "Terverifikasi: Function inquiryVa & callback handler", "95%"),
        ("2. PROTEKSI HARGA MULTI-TIER", "2.1 Proteksi skema harga 3 level (Level 1: Retail, Level 2: Grosir, Level 3: Distributor)", "Terverifikasi: Database View v_products & level_product", "100%"),
        ("2. PROTEKSI HARGA MULTI-TIER", "2.2 Sistem satuan ganda (Botol/Pcs ke Dus/Karton dengan konversi otomatis)", "Terverifikasi: Modul keranjang unit_type 1 & 2", "100%"),
        ("2. PROTEKSI HARGA MULTI-TIER", "2.3 Kalkulasi potongan promo dan diskon kupon langsung di sisi server", "Terverifikasi: Query promo aktif CAST(expired_date AS DATE)", "95%"),
        ("3. KONTROL LIMIT KREDIT", "3.1 Validasi otomatis pencegahan pemesanan melebihi limit plafon (max_credit)", "Terverifikasi: Model Mobile_api_model checkout validation", "100%"),
        ("3. KONTROL LIMIT KREDIT", "3.2 Jalur pembayaran cadangan Transfer Manual Multi-Bank dengan verifikasi kasir", "Terverifikasi: Endpoint POST /api/v1/orders/{id}/confirm-transfer", "90%")
    ],
    "96%",
    ["View database v_products mengunci skema harga multi-level secara tamper-proof.", "Dynamic VA 15 menit berhasil mencegah penumpukan pesanan pending yang menahan stok."],
    ["Risiko gangguan jaringan perbankan saat lonjakan pesanan puncak musim tanam."],
    ["Menyiapkan jalur transfer bank manual multi-bank dengan verifikasi cepat tim kasir."]
)

# ==================== SLIDE 10: BSC 8 - SISI TEKNIS ====================
create_exact_bsc_slide(
    10,
    "SISI TEKNIS (BACKEND ARCHITECTURE & RESTFUL API)",
    "Menyediakan infrastruktur backend RESTful API yang cepat (<200ms), modular (HMVC), aman, kompatibel PHP 8.x, dengan integritas transaksi MariaDB ACID.",
    {"owner": "Tim Lead Backend Engineering", "lead": "API Latency (<200ms), Error Rate (<0.01%), ACID Transaksi DB", "freq": "Bulanan", "status": "95%", "date": "31-Agu-26"},
    [
        ("1. RESTFUL API ARCHITECTURE", "1.1 Endpoint API /api/v1 terisolasi pada modul application/modules/api", "Terverifikasi: Controller Mobile.php (827 baris kode)", "100%"),
        ("1. RESTFUL API ARCHITECTURE", "1.2 Penanganan payload JSON murni dan standarisasi response error HTTP", "Terverifikasi: Function respond & error di Mobile.php", "100%"),
        ("1. RESTFUL API ARCHITECTURE", "1.3 Konfigurasi CORS header (Access-Control-Allow-Origin: *) & HTTP verbs", "Terverifikasi: Preflight OPTIONS handler di Mobile.php", "100%"),
        ("2. BASIS DATA & TRANSAKSI", "2.1 Migration scripts database (20260629_mobile_api, 20260728_account_deletion)", "Terverifikasi: Folder db/migrations/*.sql", "100%"),
        ("2. BASIS DATA & TRANSAKSI", "2.2 Penerapan transaksi ACID (trans_begin, trans_commit, trans_rollback)", "Terverifikasi: Model Mobile_api_model.php (1612 baris kode)", "100%"),
        ("2. BASIS DATA & TRANSAKSI", "2.3 Optimasi query view v_products dan pemanfaatan indeks tabel relasi", "Terverifikasi: sql_view.sql & indeks foreign key dump", "95%"),
        ("3. KINERJA & SKALABILITAS", "3.1 Kompatibilitas penuh CodeIgniter 3 HMVC pada lingkungan PHP 8.x", "Terverifikasi: Log error testing & library compatibility", "95%"),
        ("3. KINERJA & SKALABILITAS", "3.2 Keranjang belanja mobile berbasis basis data tanpa dependensi sesi web", "Terverifikasi: Tabel mobile_cart_items", "100%"),
        ("3. KINERJA & SKALABILITAS", "3.3 Cakupan pengujian unit test API otomatis", "Status: Perluasan skenario pengujian beban transaksi", "85%")
    ],
    "95%",
    ["Arsitektur API mandiri berbasis bearer token tanpa ketergantungan sesi web browser.", "Transaksi checkout terlindungi mekanisme trans_begin dan trans_commit ACID."],
    ["Kebutuhan penanganan query pada tabel historis pesanan yang terus membesar."],
    ["Menerapkan strategi pengarsipan data (archiving) transaksi di atas 2 tahun secara berkala."]
)

# ==================== SLIDE 11: BSC 9 - SISI KEAMANAN ====================
create_exact_bsc_slide(
    11,
    "SISI KEAMANAN (SECURITY & COMPLIANCE)",
    "Menjamin keamanan data kredensial pengguna, enkripsi token sesi, perlindungan hak privasi non-PII, dan kepatuhan audit regulasi global.",
    {"owner": "Tim Security & Compliance", "lead": "Zero Critical Vulnerability, Token Expiry, Audit Trail Non-PII", "freq": "Bulanan", "status": "97%", "date": "31-Agu-26"},
    [
        ("1. AUTENTIKASI & SESI", "1.1 Bearer Token SHA-256 (mobile_api_tokens) dengan kedaluwarsa 30 hari", "Terverifikasi: issue_token & user_from_token SHA-256", "100%"),
        ("1. AUTENTIKASI & SESI", "1.2 Hashing kata sandi pengguna berstandar industri BCRYPT (password_hash)", "Terverifikasi: Register & login verify di Mobile.php", "100%"),
        ("1. AUTENTIKASI & SESI", "1.3 Pencabutan token instan saat logout dan pembatasan hak akses level", "Terverifikasi: Function revoke_token & require_fields", "100%"),
        ("2. PRIVASI & NON-PII", "2.1 Endpoint penghapusan akun mandiri DELETE /api/v1/account", "Terverifikasi: App Store Guideline 5.1.1(v) compliance", "100%"),
        ("2. PRIVASI & NON-PII", "2.2 Audit trail non-PII dengan hash SHA-256 pada tabel mobile_account_deletions", "Terverifikasi: Struktur kolom email_hash CHAR(64)", "100%"),
        ("2. PRIVASI & NON-PII", "2.3 Pelepasan relasi PII pada pesanan tanpa menghapus arsip pembukuan", "Terverifikasi: Anonymization logic di delete_account", "100%"),
        ("3. OPERASIONAL & AUDIT", "3.1 Isolasi akun uji coba internal (is_internal) agar tidak mengotori omset riil", "Terverifikasi: Pemisahan data akun demo di database", "95%"),
        ("3. OPERASIONAL & AUDIT", "3.2 Sanitasi input payload JSON mencegah ancaman SQL Injection & XSS", "Terverifikasi: Input casting & query bindings", "95%"),
        ("3. OPERASIONAL & AUDIT", "3.3 Pembaruan rutin sertifikat keamanan SSL/TLS dan private key perbankan", "Terverifikasi: File key/private.pem & HTTPS headers", "90%")
    ],
    "97%",
    ["Kepatuhan penuh standar Apple App Review 5.1.1(v) dan enkripsi SHA-256 token.", "Pemisahan data audit non-PII menjamin kerahasiaan data pengguna."],
    ["Kebutuhan rotasi kunci privat enkripsi perbankan berkala sesuai kebijakan perbankan."],
    ["Menjadwalkan rotasi kunci per semester dan memonitor anomali otentikasi server."]
)

# ==================== SLIDE 12: MATRIKS TERPADU ====================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_bg(s12)

tb12 = s12.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
p12_0 = tb12.text_frame.paragraphs[0]
p12_0.text = "MATRIKS GROUNDED TERPADU: ACCOMPLISHMENT, ISSUES & RISK MITIGATION"
p12_0.font.size = Pt(16)
p12_0.font.bold = True
p12_0.font.color.rgb = C_DARK_TEXT

p12_1 = tb12.text_frame.add_paragraph()
p12_1.text = "Evaluasi Komparatif 9 Sisi Kinerja Berbasis Data Nyata Codebase dan Operasional Lapangan"
p12_1.font.size = Pt(10)
p12_1.font.color.rgb = C_MUTED_TEXT

t_shape12 = s12.shapes.add_table(10, 4, Inches(0.6), Inches(1.3), Inches(12.133), Inches(5.6))
t12 = t_shape12.table
t12.columns[0].width = Inches(2.0)
t12.columns[1].width = Inches(3.3)
t12.columns[2].width = Inches(3.3)
t12.columns[3].width = Inches(3.533)

headers12 = ["Perspektif BSC", "Accomplishment (Teruji)", "Issues & Root Cause (Tantangan)", "Next Steps & Risk Mitigation"]
for j, h in enumerate(headers12):
    cell = t12.cell(0, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_ORANGE_HEADER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

matrix_full = [
    ("1. Apps / Platform", "Lolos Apple Review 5.1.1 & Guest Browsing aktif 100%.", "Kios konvensional butuh adaptasi alur digital.", "Rilis fitur Quick Re-Order 1-klik & video panduan."),
    ("2. Customer", "Pemesanan mandiri 24/7 & live chat CS terintegrasi.", "Kios daerah terbiasa pesan via telepon pribadi sales.", "Insentif kupon promo belanja perdana via mobile."),
    ("3. Sales", "Penugasan toko binaan terkunci rapi per salesman_id.", "Kekhawatiran pemotongan komisi pesanan mobile.", "SK Direksi: komisi 100% tetap milik sales pembina."),
    ("4. Logistik", "Ongkir presisi kecamatan & quote lock 30 menit teruji.", "Tarif kurir reguler mahal untuk pupuk tonase besar.", "Aktivasi opsi armada truk internal Karisma."),
    ("5. Multiplatform", "Tri-Platform iOS Swift, Android Kotlin & Web Admin siap.", "Variasi spesifikasi ponsel Android mitra di daerah.", "Menjaga ukuran APK < 15MB & kompresi foto lokal."),
    ("6. Perusahaan", "Settlement kas cepat memotong DSO; audit faktur aman.", "Penyesuaian kebiasaan staf cabang ke sistem otomatis.", "Pelatihan SOP cabang & monitoring omset terpusat."),
    ("7. Keuangan", "BRIVA auto-settlement & margin 3 level terkunci.", "Risiko lonjakan antrean saat puncak musim tanam.", "Idempotency random external-id & fallback transfer kasir."),
    ("8. Teknis", "Backend RESTful API /api/v1 & transaksi ACID DB.", "Volume data transaksi historis semakin membesar.", "Penjadwalan archiving data berkala & optimasi indeks."),
    ("9. Keamanan", "Bearer Token SHA-256, BCRYPT, audit non-PII lulus uji.", "Kebutuhan rotasi kunci privat enkripsi berkala.", "Prosedur rotasi kunci tahunan & monitoring log server.")
]

for i, (mw, acc, isu, sol) in enumerate(matrix_full):
    row_idx = i + 1
    for j, val in enumerate([mw, acc, isu, sol]):
        cell = t12.cell(row_idx, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY_LIGHT if i % 2 == 1 else C_WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(8)
        p.font.color.rgb = C_DARK_TEXT
        if j == 0:
            p.font.bold = True

# ==================== SLIDE 13: STRATEGIC ROADMAP ====================
s13 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_bg(s13)

tb13 = s13.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
p13_0 = tb13.text_frame.paragraphs[0]
p13_0.text = "STRATEGIC ROADMAP: JADWAL & TAHAPAN IMPLEMENTASI (Q3 2026 – Q1 2027)"
p13_0.font.size = Pt(16)
p13_0.font.bold = True
p13_0.font.color.rgb = C_DARK_TEXT

p13_1 = tb13.text_frame.add_paragraph()
p13_1.text = "Rencana Kerja Bertahap Menuju Peluncuran Penuh & Skalabilitas Nasional"
p13_1.font.size = Pt(10)
p13_1.font.color.rgb = C_MUTED_TEXT

phases = [
    ("Fase 1: Peluncuran & Aktivasi", "Q3 2026 (SEGERA)", C_ORANGE_HEADER, [
        "Publikasi serentak di Apple App Store & Google Play Store.",
        "Onboarding 500+ mitra kios binaan bersama tim sales lapangan.",
        "Aktivasi promo belanja perdana via aplikasi mobile.",
        "Monitoring settlement harian BRIVA dan kecepatan kirim gudang."
    ]),
    ("Fase 2: Optimalisasi Logistik & Komisi", "Q4 2026", C_GREEN_STATUS, [
        "Integrasi opsi armada truk internal untuk muatan tonase besar.",
        "Peluncuran program poin loyalitas dan reward toko aktif.",
        "Pemberlakuan penuh skema komisi digital bagi sales lapangan.",
        "Otomasi promo diskon musiman cuci gudang deadstock."
    ]),
    ("Fase 3: Konsolidasi & Skalabilitas", "Q1 2027", C_CYAN_ACCENT, [
        "Konsolidasi peramalan kebutuhan stok berbasis siklus tanam.",
        "Perluasan area binaan cabang ke sentra pertanian baru.",
        "Peningkatan kapasitas server dan fitur kemitraan lanjutan.",
        "Integrasi sistem konsultasi teknis budidaya tanaman."
    ])
]

for idx, (p_title, p_time, p_col, p_items) in enumerate(phases):
    px = Inches(0.6 + idx * 4.1)
    p_card = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(1.4), Inches(3.9), Inches(5.3))
    p_card.fill.solid()
    p_card.fill.fore_color.rgb = C_NAVY_LIGHT
    p_card.line.color.rgb = C_BORDER
    
    top_strip = s13.shapes.add_shape(MSO_SHAPE.RECTANGLE, px, Inches(1.4), Inches(3.9), Inches(0.1))
    top_strip.fill.solid()
    top_strip.fill.fore_color.rgb = p_col
    top_strip.line.color.rgb = p_col
    
    p_tb = s13.shapes.add_textbox(px + Inches(0.2), Inches(1.6), Inches(3.5), Inches(4.9))
    p_tf = p_tb.text_frame
    p_tf.word_wrap = True
    
    pp0 = p_tf.paragraphs[0]
    pp0.text = p_title
    pp0.font.size = Pt(12)
    pp0.font.bold = True
    pp0.font.color.rgb = C_DARK_TEXT
    
    pp1 = p_tf.add_paragraph()
    pp1.text = p_time
    pp1.font.size = Pt(10)
    pp1.font.bold = True
    pp1.font.color.rgb = p_col
    pp1.space_before = Pt(2)
    
    for item in p_items:
        pp = p_tf.add_paragraph()
        pp.text = f"• {item}"
        pp.font.size = Pt(9)
        pp.font.color.rgb = C_DARK_TEXT
        pp.space_before = Pt(6)

# ==================== SLIDE 14: KESIMPULAN EKSEKUTIF ====================
s14 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_bg(s14)

tb14 = s14.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
p14_0 = tb14.text_frame.paragraphs[0]
p14_0.text = "KESIMPULAN EKSEKUTIF & PERSETUJUAN DIREKSI"
p14_0.font.size = Pt(16)
p14_0.font.bold = True
p14_0.font.color.rgb = C_DARK_TEXT

p14_1 = tb14.text_frame.add_paragraph()
p14_1.text = "Persetujuan Peluncuran Komersial Ekosistem Karisma Online"
p14_1.font.size = Pt(10)
p14_1.font.color.rgb = C_MUTED_TEXT

# Left Card: Ringkasan Evaluasi
c_l = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.9), Inches(4.5))
c_l.fill.solid()
c_l.fill.fore_color.rgb = C_NAVY_LIGHT
c_l.line.color.rgb = C_BORDER

tb_l = s14.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.1))
tf_l = tb_l.text_frame
tf_l.word_wrap = True

pl0 = tf_l.paragraphs[0]
pl0.text = "📌 Ringkasan Evaluasi 9 Sisi Kesiapan"
pl0.font.size = Pt(13)
pl0.font.bold = True
pl0.font.color.rgb = C_DARK_TEXT

eval_items = [
    "Rata-Rata Kesiapan 94.2%: Ekosistem Karisma Online telah mencapai kesiapan komersial penuh dan lulus seluruh audit regulasi.",
    "Kepatuhan Standar Industri: Memenuhi 100% persyaratan Apple App Store, Google Play Store, dan perbankan Bank BRI.",
    "Efisiensi Operasional Terbukti: Pemangkasan waktu settlement kas dari jam ke detik serta akurasi ongkos kirim tingkat kecamatan."
]
for item in eval_items:
    pl = tf_l.add_paragraph()
    pl.text = f"• {item}"
    pl.font.size = Pt(10)
    pl.font.color.rgb = C_DARK_TEXT
    pl.space_before = Pt(8)

# Right Card: Rekomendasi
c_r = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.9), Inches(4.5))
c_r.fill.solid()
c_r.fill.fore_color.rgb = C_NAVY_LIGHT
c_r.line.color.rgb = C_BORDER

tb_r = s14.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.1))
tf_r = tb_r.text_frame
tf_r.word_wrap = True

pr0 = tf_r.paragraphs[0]
pr0.text = "✅ Rekomendasi Keputusan Dewan Direksi"
pr0.font.size = Pt(13)
pr0.font.bold = True
pr0.font.color.rgb = C_GREEN_STATUS

rec_items = [
    "1. Persetujuan Peluncuran (Go-Live Sign-Off): Memberikan persetujuan resmi peluncuran publik aplikasi mobile Karisma Online.",
    "2. Kebijakan Komisi Digital Sales: Mengesahkan SK Direksi terkait alokasi 100% komisi pesanan digital untuk salesman pembina kios.",
    "3. Alokasi Armada Logistik: Menyetujui skema pengiriman armada truk internal untuk pesanan pupuk tonase besar."
]
for item in rec_items:
    pr = tf_r.add_paragraph()
    pr.text = f"• {item}"
    pr.font.size = Pt(10)
    pr.font.color.rgb = C_DARK_TEXT
    pr.space_before = Pt(8)

# Bottom Bar
bot_b = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.2), Inches(12.133), Inches(0.7))
bot_b.fill.solid()
bot_b.fill.fore_color.rgb = C_NAVY_DARK
bot_b.line.color.rgb = C_NAVY_DARK

bot_tb = s14.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.5))
bp0 = bot_tb.text_frame.paragraphs[0]
bp0.text = "Dokumen ini siap disahkan sebagai acuan kerja resmi PT. Karisma Indoagro Universal. Sesi Diskusi & Tanya Jawab (Q&A) Dipersilakan."
bp0.font.size = Pt(10.5)
bp0.font.bold = True
bp0.font.color.rgb = C_WHITE

output_path = "/Applications/XAMPP/xamppfiles/htdocs/kiustore/docs/nw_presentasi_karisma_online.pptx"
prs.save(output_path)
print(f"SUCCESS: Generated 9 BSC presentation deck at {output_path}")
