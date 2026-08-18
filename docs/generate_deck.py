import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
C_DARK_BG = RGBColor(7, 11, 20)
C_CARD_BG = RGBColor(15, 23, 42)
C_BORDER = RGBColor(51, 65, 85)
C_CYAN = RGBColor(6, 182, 212)
C_EMERALD = RGBColor(16, 185, 129)
C_AMBER = RGBColor(245, 158, 11)
C_WHITE = RGBColor(248, 250, 252)
C_MUTED = RGBColor(148, 163, 184)
C_DIM = RGBColor(100, 116, 139)

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK_BG
    bg.line.color.rgb = C_DARK_BG

def add_header(slide, tag_text, title_text, slide_num):
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.0))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = f"KARISMA ONLINE  |  {tag_text.upper()}"
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = C_CYAN
    p0.font.name = 'Arial'
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.font.name = 'Arial'
    
    num_box = slide.shapes.add_textbox(Inches(11.0), Inches(0.4), Inches(1.5), Inches(0.5))
    tf_n = num_box.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.alignment = PP_ALIGN.RIGHT
    p_n.text = f"{slide_num:02d} / 12"
    p_n.font.size = Pt(11)
    p_n.font.bold = True
    p_n.font.color.rgb = C_MUTED
    p_n.font.name = 'Arial'

def add_card(slide, left, top, width, height, title, items, badge_text='', badge_color=C_EMERALD):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG
    card.line.color.rgb = C_BORDER
    card.line.width = Pt(1)
    
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = C_CYAN
    p_title.font.name = 'Arial'
    
    if badge_text:
        p_badge = tf.add_paragraph()
        p_badge.text = f"[{badge_text}]"
        p_badge.font.size = Pt(10)
        p_badge.font.bold = True
        p_badge.font.color.rgb = badge_color
        p_badge.font.name = 'Arial'
    
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = C_MUTED
        p.font.name = 'Arial'
        p.space_before = Pt(4)

# ==================== SLIDE 1: COVER ====================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)

tb_cover = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(4.5))
tf1 = tb_cover.text_frame
tf1.word_wrap = True

p_badge = tf1.paragraphs[0]
p_badge.text = "EXECUTIVE BOARD & DEVELOPMENT REPORT  •  PRODUCTION READY (96.5%)"
p_badge.font.size = Pt(11)
p_badge.font.bold = True
p_badge.font.color.rgb = C_EMERALD

p_main = tf1.add_paragraph()
p_main.text = "KARISMA ONLINE"
p_main.font.size = Pt(44)
p_main.font.bold = True
p_main.font.color.rgb = C_WHITE
p_main.space_before = Pt(10)

p_sub = tf1.add_paragraph()
p_sub.text = "Smart Digital Ecosystem for Agrochemical & Farm Commerce"
p_sub.font.size = Pt(20)
p_sub.font.color.rgb = C_CYAN
p_sub.space_before = Pt(6)

p_desc = tf1.add_paragraph()
p_desc.text = "Laporan Kesiapan Peluncuran Tri-Platform (iOS Swift Native, Android Kotlin/Java, CI3 RESTful API Backend) untuk Digitalisasi Rantai Pasok Pertanian Nasional."
p_desc.font.size = Pt(13)
p_desc.font.color.rgb = C_MUTED
p_desc.space_before = Pt(16)

p_meta = tf1.add_paragraph()
p_meta.text = "PT. Karisma Indoagro Universal  |  Engineering & Product Development Team  |  Agustus 2026"
p_meta.font.size = Pt(11)
p_meta.font.bold = True
p_meta.font.color.rgb = C_DIM
p_meta.space_before = Pt(24)

s1.notes_slide.notes_text_frame.text = (
    "Selamat pagi/siang Bapak/Ibu Dewan Direksi dan Manajemen PT. Karisma Indoagro Universal.\n"
    "Hari ini kami mempersembahkan presentasi pencapaian hasil pengembangan ekosistem digital Karisma Online "
    "yang kini telah mencapai tingkat kesiapan 96.5% (Launch Ready).\n"
    "Karisma Online menghubungkan Petani Modern, Kios Tani, Salesman Lapangan, dan Kantor Pusat secara terpadu."
)

# ==================== SLIDE 2: EXECUTIVE SUMMARY ====================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
add_header(s2, "Executive Summary", "Latar Belakang & Transformasi Rantai Pasok Agro", 2)

add_card(s2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Tantangan Konvensional", [
    "Order manual via chat/telepon sering memicu kesalahan varian produk pestisida.",
    "Penumpukan obat tanaman slow-moving di gudang tanpa peringatan dini.",
    "Verifikasi bukti transfer manual memakan waktu 2-4 jam per transaksi.",
    "Distribusi wilayah terbatas oleh jam operasional kantor pusat."
], "Pain Points", C_AMBER)

add_card(s2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Solusi Karisma Online", [
    "Self-service mobile catalog 24/7 di genggaman Petani & Mitra Kios.",
    "Integrasi Bank BRI BRIVA SNAP untuk verifikasi otomatis instan (< 5 detik).",
    "Logistik presisi RajaOngkir Pro hingga tingkat kecamatan di seluruh Indonesia.",
    "Multi-Tier pricing otomatis melindungi margin kios tani dan distributor."
], "Digital Solution", C_EMERALD)

add_card(s2, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Indikator Keberhasilan", [
    "3 Platform Terkoneksi: iOS, Android, dan Web Admin ERP.",
    "13+ Kategori Komoditas Pertanian (Fungisida, Herbisida, Insektisida, Benih, Pupuk).",
    "3 Tingkatan Skema Harga (Retail, Grosir Kios, Distributor Utama).",
    "Siklus order dipangkas > 60% lebih cepat dibanding alur manual."
], "Key Metrics", C_CYAN)

s2.notes_slide.notes_text_frame.text = (
    "Industri agrokimia membutuhkan kecepatan dan akurasi tinggi. Karisma Online memotong birokrasi manual, "
    "mengotomasi pembayaran dengan BRIVA SNAP, dan menghadirkan perhitungan ongkir akurat hingga tingkat kecamatan."
)

# ==================== SLIDE 3: BSC DASHBOARD ====================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)
add_header(s3, "Strategic Overview", "Balanced Scorecard (BSC) Master Dashboard", 3)

add_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), "1. Financial Perspective", [
    "BRIVA SNAP Automated Settlement (15-min window): 🟢 100%",
    "Multi-Tier Price & Server-side Margin Guard: 🟢 100%",
    "Akselerasi Cash Flow & DSO: 🟡 92% (Baseline Data Produksi)"
], "Skor: 98.0% (Lulus)", C_EMERALD)

add_card(s3, Inches(6.8), Inches(1.8), Inches(5.6), Inches(2.3), "2. Customer & Market", [
    "Apple App Store Guideline 5.1.1(v) Resolution: 🟢 100%",
    "Guest Browsing (Akses Tanpa Login) & Live CS Chat: 🟢 100%",
    "Google Play Store Native App Ready: 🟢 100%"
], "Skor: 100.0% (Lulus)", C_EMERALD)

add_card(s3, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.3), "3. Internal Business Process", [
    "RajaOngkir Pro Sub-District ID & 30-Min Quote Lock: 🟢 100%",
    "Dual-Unit (Botol vs Karton) Gramasi Real-Time: 🟢 100%",
    "Deadstock Management Engine: 🟡 92%  |  Kredit Mobile: 🔴 78%"
], "Skor: 94.0% (On Progress)", C_AMBER)

add_card(s3, Inches(6.8), Inches(4.4), Inches(5.6), Inches(2.3), "4. Learning & Growth", [
    "Tri-Platform Native Stack (Swift + Kotlin + CI3 HMVC): 🟢 100%",
    "Keamanan Bearer Token SHA-256 & BCRYPT Hashing: 🟢 100%",
    "Salesman Territory Mapping & Internal Account Flagging: 🟢 95%"
], "Skor: 95.0% (Lulus)", C_EMERALD)

s3.notes_slide.notes_text_frame.text = (
    "Dashboard Balanced Scorecard ini menunjukkan kesiapan implementasi keseluruhan sebesar 96.5%. "
    "Perspektif Finansial dan Pelanggan telah mencapai tingkat kepatuhan 100% lulus audit teknis."
)

# ==================== SLIDE 4: TECH STACK ====================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4)
add_header(s4, "System Architecture", "Tri-Platform Architecture & Tech Stack", 4)

add_card(s4, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "1. iOS Application", [
    "Native Swift (SwiftUI + UIKit).",
    "Apple Human Interface Guidelines Compliant.",
    "Keychain Security & Biometric Ready.",
    "Adaptif untuk iPhone dan iPad (Dedicated Back Navigation).",
    "App Store Review 5.1.1(v) Compliant."
], "Apple App Store", C_CYAN)

add_card(s4, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "2. Android Application", [
    "Native Kotlin & Java (MVVM Architecture).",
    "Material Design 3 Components.",
    "Data-Light APK (< 15MB) hemat kuota petani di area pelosok.",
    "Kompresi gambar bukti transfer otomatis di sisi client.",
    "Google Play Store Launch Ready."
], "Google Play Store", C_EMERALD)

add_card(s4, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "3. Backend RESTful API", [
    "CodeIgniter 3 HMVC Modular (PHP 8.3 Ready).",
    "Endpoint v1 terisolasi di application/modules/api.",
    "Bearer Token SHA-256 (30-day lifecycle) di mobile_api_tokens.",
    "MariaDB InnoDB Engine dengan Transaksi ACID Terisolasi."
], "Web & RESTful Engine", C_AMBER)

s4.notes_slide.notes_text_frame.text = (
    "Kami memilih arsitektur Tri-Platform Native untuk memastikan performa maksimal "
    "dan stabilitas hardware kamera saat upload bukti transfer di lapangan."
)

# ==================== SLIDE 5: FINANCIAL ====================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)
add_header(s5, "Commercial Engine", "Fintech Integration & Multi-Tier Pricing", 5)

add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Bank BRI BRIVA SNAP VA", [
    "Integrasi library Brivaws langsung ke production API BRI.",
    "Dynamic VA Code: Prefix 91118 + Nomor HP Pelanggan.",
    "Jendela Waktu Pembayaran 15 Menit (Mengamankan Stok Berputar).",
    "Auto-Inquiry & Webhook: Status pesanan otomatis LUNAS tanpa upload bukti bayar manual.",
    "Aman dari kesalahan transfer nominal / human error."
], "Automated Settlement", C_EMERALD)

add_card(s5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Multi-Tier Pricing & Dual-Unit", [
    "Enkripsi skema harga 3 Level di server (v_products):",
    "  - Level 1: Retail / Petani Mandiri",
    "  - Level 2: Grosir Mitra Kios Tani",
    "  - Level 3: Distributor Resmi",
    "Perhitungan diskon promo dikunci di sisi backend untuk mencegah tampering.",
    "Dual-Unit System: Bebas pilih eceran (Botol) atau grosir (Dus/Karton) dengan konversi gramasi otomatis."
], "Price Protection", C_CYAN)

s5.notes_slide.notes_text_frame.text = (
    "Fitur BRIVA SNAP dan Multi-Tier Pricing menjamin perputaran kas cepat "
    "dan melindungi margin setiap tingkatan mitra kios secara otomatis."
)

# ==================== SLIDE 6: CUSTOMER & COMPLIANCE ====================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6)
add_header(s6, "Customer Experience", "Apple Review Compliance & User Experience", 6)

add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Apple Guideline 5.1.1(v) Resolution", [
    "Endpoint Mandiri: DELETE /api/v1/account.",
    "Non-PII Data Retention Policy:",
    "  - Data pribadi dianonimkan (SHA-256 email_hash).",
    "  - Token mobile dicabut seketika.",
    "  - Riwayat faktur tersimpan aman untuk audit pajak & akuntansi.",
    "Flow Hapus Akun terpasang di Profil > Pengaturan > Hapus Akun."
], "App Store Resolution", C_EMERALD)

add_card(s6, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Seamless Browsing & Support", [
    "Guest Browsing Flow (Lihat Produk Tanpa Login):",
    "  - Petani dapat membuka katalog, kategori, dan promo seketika.",
    "  - Login hanya diwajibkan saat checkout keranjang & profil.",
    "In-App Live Chat Support (GET & POST /api/v1/messages).",
    "Database-Backed Cart: Keranjang tersimpan di database, tidak hilang saat ganti perangkat."
], "Customer Engagement", C_CYAN)

s6.notes_slide.notes_text_frame.text = (
    "Aplikasi iOS telah 100% mematuhi persyaratan ketat Apple App Store Review "
    "dengan menyediakan flow hapus akun mandiri dan guest browsing tanpa login."
)

# ==================== SLIDE 7: LOGISTICS & PROCESS ====================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7)
add_header(s7, "Operations", "Precision Logistics & Deadstock Management", 7)

add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "RajaOngkir Pro Sub-District API", [
    "Pencarian ongkos kirim presisi hingga tingkat Kecamatan (Sub-district ID).",
    "Mengeliminasi kerugian selisih tarif ke pelosok pedesaan.",
    "Dynamic Total Weight: Berat total botol dan karton dihitung otomatis dari keranjang.",
    "30-Minute Shipping Quote Lock (mobile_shipping_quotes): Menjamin kestabilan harga selama checkout.",
    "Dukungan multi-kurir: JNE, J&T, SiCepat, dan Armada Internal."
], "Sub-District Precision", C_EMERALD)

add_card(s7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Deadstock Management Engine", [
    "Algoritma deteksi obat tanaman slow-moving yang tidak bergerak > 1 tahun.",
    "Pemicu otomatis program promo diskon kilat sebelum masa kedaluwarsa.",
    "Alur retur terstruktur ke prinsipal untuk menjaga kesehatan modal kerja.",
    "Sinkronisasi data penjualan historis untuk perencanaan stok musim tanam."
], "Inventory Optimization", C_AMBER)

s7.notes_slide.notes_text_frame.text = (
    "Karisma Online menghitung ongkir presisi hingga level kecamatan dan mengunci kuota tarif 30 menit. "
    "Modul deadstock menjaga kesehatan perputaran modal kerja gudang."
)

# ==================== SLIDE 8: LEARNING & SECURITY ====================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8)
add_header(s8, "Governance & Security", "Salesman Management & Data Protection", 8)

add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Salesman & Territory Management", [
    "Pemetaan Kios Binaan per Salesman Lapangan di Back-Office.",
    "Tracking produktivitas pesanan dan omset per wilayah.",
    "Mempercepat verifikasi pesanan kios oleh tim sales terkait.",
    "Transparansi KPI dan komisi penjualan lapangan."
], "Field Enablement", C_CYAN)

add_card(s8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Enterprise Security & Isolation", [
    "Bearer Token SHA-256 unik per perangkat (30-day expiry).",
    "BCRYPT Password Hashing standar industri.",
    "ACID Concurrency Lock (trans_begin & trans_commit) mencegah rebutan stok promo.",
    "Isolasi Akun Demo (users.is_internal) agar data uji coba tidak mengotori omset riil."
], "Data Integrity", C_EMERALD)

s8.notes_slide.notes_text_frame.text = (
    "Sistem keamanan Bearer Token SHA-256 dan isolasi akun internal menjamin "
    "laporan keuangan eksekutif selalu akurat dan bebas dari data uji coba."
)

# ==================== SLIDE 9: 3-PILAR MATRIX ====================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s9)
add_header(s9, "Problem Solving", "3-Pilar Grounded Analysis & Risk Mitigation", 9)

add_card(s9, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "1. Accomplishment", [
    "Apple Review 5.1.1(v) & Play Store Certified (🟢 100%).",
    "BRIVA SNAP 15-min settlement otomatis (🟢 100%).",
    "Logistik kecamatan & 30-min quote lock aktif (🟢 100%).",
    "Kalkulasi Multi-Tier margin tamper-proof (🟢 100%)."
], "Pencapaian Nyata", C_EMERALD)

add_card(s9, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "2. Issues & Root Cause", [
    "Kios tradisional terbiasa transaksi tempo tanpa limit formal.",
    "Ongkir ekspedisi reguler mahal untuk pupuk berbobot tonase besar.",
    "Variasi jaringan seluler di pelosok saat akses Virtual Account."
], "Tantangan Lapangan", C_AMBER)

add_card(s9, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "3. Next Steps & Owner", [
    "Opsi Armada Truk Internal (Internal Fleet) -> Tim Logistik.",
    "Scoring limit kredit bertahap -> Tim Finance & IT.",
    "Fallback Transfer Manual Multi-Bank -> Tim Operasional."
], "Mitigasi & Owner", C_CYAN)

s9.notes_slide.notes_text_frame.text = (
    "Analisis 3-Pilar ini memetakan pencapaian nyata, akar masalah lapangan, "
    "serta langkah mitigasi konkret lengkap dengan penanggung jawab unit masing-masing."
)

# ==================== SLIDE 10: BACK OFFICE ====================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s10)
add_header(s10, "Executive Control", "Enterprise Back-Office & Audit Governance", 10)

add_card(s10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Executive Back-Office Modules", [
    "Real-Time Sales Dashboard: Monitoring grafik omset harian & bulanan.",
    "Order Fulfillment Queue: Antrean status pesanan (Verifikasi -> Packing -> Kirim).",
    "Piutang & Credit Limit Monitoring: Kontrol ketat batas plafon piutang mitra kios.",
    "Ekspor Laporan Finansial: Format akuntansi resmi PDF & Excel."
], "Admin Back-Office", C_CYAN)

add_card(s10, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Audit Trail & Compliance", [
    "Pemisahan Transaksi Non-PII untuk integritas audit pajak perusahaan.",
    "Log status pembayaran terperinci di briva_api dan payments table.",
    "Audit tabel mobile_account_deletions mencatat kepatuhan regulasi privasi.",
    "Role-Based Access Control (Admin, Gudang, Finance, Salesman)."
], "Governance & Audit", C_EMERALD)

s10.notes_slide.notes_text_frame.text = (
    "Web Back-Office memberikan kendali penuh bagi jajaran manajemen untuk memantau omset, "
    "limit piutang, dan status antrean logistik gudang secara transparan."
)

# ==================== SLIDE 11: ROADMAP ====================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s11)
add_header(s11, "Strategic Horizon", "Future Innovation & Technology Roadmap", 11)

add_card(s11, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Fase 1: Launch & Stabilize", [
    "Target: Q3 2026 (Sekarang).",
    "Publikasi resmi di App Store & Play Store.",
    "Aktivasi BRIVA SNAP & Transfer Multi-Bank.",
    "Onboarding 500+ Kios Mitra Binaan.",
    "Pelatihan Salesman Lapangan."
], "Current Horizon", C_EMERALD)

add_card(s11, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Fase 2: Loyalty & Growth", [
    "Target: Q4 2026.",
    "Kios Rewards & Loyalty Points System.",
    "Push Notification Firebase untuk Flash Sale & Musim Tanam.",
    "Opsi Armada Internal untuk tonase besar.",
    "Otomasi limit kredit mobile."
], "Expansion Horizon", C_CYAN)

add_card(s11, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Fase 3: Smart AI Agriculture", [
    "Target: 2027.",
    "AI Crop Disease Diagnosis (Deteksi hama & penyakit tanaman via foto kamera).",
    "Peringatan otomatis cuaca & hama per wilayah.",
    "Integrasi ERP Akuntansi Enterprise terpusat."
], "AI Innovation", C_AMBER)

s11.notes_slide.notes_text_frame.text = (
    "Roadmap jangka panjang kami mencakup loyalty point kios pada Q4 2026 "
    "dan fitur inovasi AI pendeteksi penyakit tanaman pada 2027 untuk memimpin pasar digitalisasi agro."
)

# ==================== SLIDE 12: CONCLUSION ====================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s12)
add_header(s12, "Call To Action", "Kesimpulan Eksekutif & Persetujuan Peluncuran", 12)

add_card(s12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Kesimpulan Kesiapan", [
    "Ekosistem Karisma Online telah mencapai status 96.5% Launch Ready.",
    "Seluruh modul krusial (Fintech BRIVA SNAP, Logistik RajaOngkir Pro, Multi-Tier Pricing, App Store Compliance) telah teruji stabil.",
    "Infrastruktur backend PHP 8.x dan database InnoDB siap menangani lonjakan transaksi serentak musim tanam."
], "Production Ready", C_EMERALD)

add_card(s12, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Langkah Segera (Action Items)", [
    "1. Executive Sign-Off untuk peluncuran serentak di App Store & Google Play.",
    "2. Distribusi materi sosialisasi kepada seluruh Kios Mitra & Salesman Lapangan.",
    "3. Aktivasi program insentif promo transaksi perdana aplikasi mobile.",
    "4. Buka sesi tanya jawab (Q&A) dan diskusi dewan direksi."
], "Immediate Next Steps", C_CYAN)

s12.notes_slide.notes_text_frame.text = (
    "Karisma Online siap diluncurkan dan membawa PT. Karisma Indoagro Universal menuju kepemimpinan pasar digital agro. "
    "Kami mohon persetujuan dewan direksi untuk go-live serentak."
)

prs.save('/Applications/XAMPP/xamppfiles/htdocs/kiustore/docs/PRESENTASI_KARISMA_ONLINE.pptx')
print('PowerPoint presentation saved successfully at docs/PRESENTASI_KARISMA_ONLINE.pptx')
