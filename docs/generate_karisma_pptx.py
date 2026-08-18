#!/usr/bin/env python3
"""
Generator Presentasi PowerPoint (PPTX) BSC Karisma Online
Dilengkapi Bukti Eviden Faktual (Audit-Ready) untuk Memimpin dengan Data.
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(DOCS_DIR, "BSC_Karisma_Online.pptx")

# Color Palette: Modern Obsidian Navy & Vibrant Strategic Accents
C_BG_DARK = RGBColor(10, 15, 29)         # #0A0F1D
C_CARD_BG = RGBColor(22, 32, 50)         # #162032
C_CARD_BORDER = RGBColor(45, 60, 85)     # #2D3C55
C_GOLD_HEADER = RGBColor(245, 158, 11)   # #F59E0B
C_CYAN_ACCENT = RGBColor(6, 182, 212)    # #06B6D4
C_BLUE_LIGHT = RGBColor(56, 189, 248)    # #38BDF8
C_GREEN_STATUS = RGBColor(16, 185, 129)  # #10B981
C_YELLOW_STATUS = RGBColor(234, 179, 8)  # #EAB308
C_RED_STATUS = RGBColor(239, 68, 68)     # #EF4444
C_WHITE = RGBColor(255, 255, 255)
C_SLATE_200 = RGBColor(226, 232, 240)
C_SLATE_400 = RGBColor(148, 163, 184)
C_SLATE_700 = RGBColor(51, 65, 85)

def create_deck():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG_DARK
        bg.line.color.rgb = C_BG_DARK
        return bg

    def add_slide_header(slide, title_text, category_text="BALANCED SCORECARD (BSC) // MEMIMPIN DENGAN DATA"):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(9.5)
        p_cat.font.bold = True
        p_cat.font.color.rgb = C_CYAN_ACCENT

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = C_WHITE
        p_title.space_before = Pt(3)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = C_CARD_BORDER
        line.line.color.rgb = C_CARD_BORDER

    def add_card(slide, left, top, width, height):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_CARD_BORDER
        card.line.width = Pt(1)
        return card

    # ==========================================================================
    # SLIDE 1: COVER SLIDE
    # ==========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1)

    cover_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(0.9), Inches(10.933), Inches(4.0))
    cover_card.fill.solid()
    cover_card.fill.fore_color.rgb = C_CARD_BG
    cover_card.line.color.rgb = C_CARD_BORDER

    tb_cov = s1.shapes.add_textbox(Inches(1.6), Inches(1.2), Inches(10.133), Inches(3.4))
    tf_cov = tb_cov.text_frame
    tf_cov.word_wrap = True

    p0 = tf_cov.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    p0.text = "EXECUTIVE STRATEGIC SCORECARD // SINGLE SOURCE OF TRUTH & BUKTI EVIDEN"
    p0.font.size = Pt(10.5)
    p0.font.bold = True
    p0.font.color.rgb = C_CYAN_ACCENT

    p1 = tf_cov.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.text = "KARISMA ONLINE"
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.space_before = Pt(6)

    p2 = tf_cov.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Balanced Scorecard (BSC) & Verifikasi Bukti Eviden Kesiapan Sistem"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = C_GOLD_HEADER
    p2.space_before = Pt(4)

    p3 = tf_cov.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = "Multi-Platform • Transaksi Customer • Payment Gateway • UAT Minimal 3 Kios • SLA 30 Hari Bebas Bug • Dokumentasi & SOP"
    p3.font.size = Pt(11)
    p3.font.color.rgb = C_SLATE_400
    p3.space_before = Pt(10)

    p4 = tf_cov.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.text = "Tim Digital & IT Karisma  |  Target Go-Live Q3 2026  |  Lead With Data & Audit-Ready"
    p4.font.size = Pt(10)
    p4.font.bold = True
    p4.font.color.rgb = C_BLUE_LIGHT
    p4.space_before = Pt(14)

    hero_kpis = [
        ("89.1%", "Overall Readiness Status", "Kategori: On Progress (Dalam Pengawalan)", C_YELLOW_STATUS),
        ("7 Must Wins", "Inisiatif Kunci Terpadu", "28 Sub-Inisiatif + Bukti Eviden", C_CYAN_ACCENT),
        ("0 Critical Bug", "Target Pasca-Live (D+30)", "SLA Respon <15 Mnt, Fix <4 Jam", C_GREEN_STATUS),
        ("Min. 3 Kios", "Target Verifikasi UAT", "Syarat Mutlak Gate Review Rilis", C_BLUE_LIGHT),
    ]

    for i, (val, title, sub, col) in enumerate(hero_kpis):
        k_left = Inches(1.2 + i * 2.78)
        k_top = Inches(5.2)
        add_card(s1, k_left, k_top, Inches(2.6), Inches(1.6))

        tb_k = s1.shapes.add_textbox(k_left + Inches(0.12), k_top + Inches(0.15), Inches(2.36), Inches(1.3))
        tf_k = tb_k.text_frame
        tf_k.word_wrap = True

        p_v = tf_k.paragraphs[0]
        p_v.alignment = PP_ALIGN.CENTER
        p_v.text = val
        p_v.font.size = Pt(22)
        p_v.font.bold = True
        p_v.font.color.rgb = col

        p_t = tf_k.add_paragraph()
        p_t.alignment = PP_ALIGN.CENTER
        p_t.text = title
        p_t.font.size = Pt(10)
        p_t.font.bold = True
        p_t.font.color.rgb = C_WHITE
        p_t.space_before = Pt(2)

        p_s = tf_k.add_paragraph()
        p_s.alignment = PP_ALIGN.CENTER
        p_s.text = sub
        p_s.font.size = Pt(8)
        p_s.font.color.rgb = C_SLATE_400
        p_s.space_before = Pt(2)

    # ==========================================================================
    # SLIDE 2: SCORECARD MASTER VIEW + AUDIT EVIDEN SUMMARY
    # ==========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2)
    add_slide_header(s2, "Executive Scorecard: 7 Must Wins & Bukti Eviden Sistem")

    meta_bg = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.733), Inches(0.85))
    meta_bg.fill.solid()
    meta_bg.fill.fore_color.rgb = C_CARD_BG
    meta_bg.line.color.rgb = C_CARD_BORDER

    tb_m = s2.shapes.add_textbox(Inches(0.95), Inches(1.5), Inches(11.433), Inches(0.75))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True

    p_m1 = tf_m.paragraphs[0]
    p_m1.text = "Objective Description: Kesiapan & Stabilitas Sistem Karisma Online (Multi-Platform, Transaksi Customer, Payment Gateway, UAT 3 Kios, Bebas Bug 30 Hari, Dokumentasi & SOP)"
    p_m1.font.size = Pt(9)
    p_m1.font.bold = True
    p_m1.font.color.rgb = C_SLATE_200

    p_m2 = tf_m.add_paragraph()
    p_m2.text = "Owner: Tim IT & Digital Karisma  |  Measure Lead: Project & QA Lead  |  Frequency: Mingguan  |  Overall Status: 89% (On Progress)  |  Tanggal: 10-Agu-26"
    p_m2.font.size = Pt(8.5)
    p_m2.font.color.rgb = C_CYAN_ACCENT
    p_m2.space_before = Pt(4)

    rows, cols = 8, 4
    t_left, t_top, t_w, t_h = Inches(0.8), Inches(2.4), Inches(11.733), Inches(4.3)
    table_shape = s2.shapes.add_table(rows, cols, t_left, t_top, t_w, t_h)
    table = table_shape.table
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(6.733)
    table.columns[3].width = Inches(1.4)

    t_headers = ["No", "Must Win", "Key Initiatives & Bukti Eviden (Single Source of Truth)", "Status"]
    for c_idx, h_text in enumerate(t_headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_GOLD_HEADER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    scorecard_summary = [
        ("1", "MULTI PLATFORM", "Validasi lintas platform. Eviden: Kontrak API `docs/MOBILE_API.md`, Controller `Mobile.php`, Migrasi `20260629_mobile_api.sql`.", "92%", C_YELLOW_STATUS),
        ("2", "CUSTOMER BISA TRANSAKSI STABIL", "Validasi Order → Bayar → Konfirmasi. Eviden: `Mobile.php::checkout()`, DB `orders` & `order_items`, `curl_response_log.txt`.", "91%", C_YELLOW_STATUS),
        ("3", "INTEGRASI PAYMENT GATEWAY STABIL", "Validasi pembayaran real-time. Eviden: Controller `briva_list_function.php`, `signature_log.txt`, Webhook callback BRIVA.", "93%", C_YELLOW_STATUS),
        ("4", "UAT DENGAN MINIMAL 3 KIOS", "Uji di 3 kios pilot nyata. Eviden: Dokumen Skenario `SKENARIO-UAT-KIU-2026`, Berita Acara UAT, UAT Issue Tracker.", "85%", C_RED_STATUS),
        ("5", "TIDAK ADA BUG CRITICAL 30 HARI LIVE", "Stabilitas D+30 pasca-live. Eviden: System Log `application/logs/`, Kalender Rilis Go-Live, On-Call 24/7 Matrix.", "85%", C_RED_STATUS),
        ("6", "DOKUMENTASI", "Standarisasi deliverable. Eviden: Master Checklist Dokumen `docs/`, User Manual Mobile/Kasir, API Contract `MOBILE_API.md`.", "90%", C_YELLOW_STATUS),
        ("7", "SOP TRANSAKSI KIOS", "SOP operasional kasir riil. Eviden: Dokumen `SOP/IT-KIU/2026/001`, Flowchart Kasir Kios, Berita Acara Training.", "88%", C_RED_STATUS),
    ]

    for r_idx, (no, mw, ki, stat, st_col) in enumerate(scorecard_summary, 1):
        c0 = table.cell(r_idx, 0)
        c0.fill.solid()
        c0.fill.fore_color.rgb = C_CARD_BG
        c0.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = c0.text_frame.paragraphs[0]
        p0.text = no
        p0.font.size = Pt(9.5)
        p0.font.bold = True
        p0.font.color.rgb = C_WHITE
        p0.alignment = PP_ALIGN.CENTER

        c1 = table.cell(r_idx, 1)
        c1.fill.solid()
        c1.fill.fore_color.rgb = C_CARD_BG
        c1.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = c1.text_frame.paragraphs[0]
        p1.text = mw
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = C_WHITE

        c2 = table.cell(r_idx, 2)
        c2.fill.solid()
        c2.fill.fore_color.rgb = C_CARD_BG
        c2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = c2.text_frame.paragraphs[0]
        p2.text = ki
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = C_SLATE_200

        c3 = table.cell(r_idx, 3)
        c3.fill.solid()
        c3.fill.fore_color.rgb = C_CARD_BG
        c3.vertical_anchor = MSO_ANCHOR.MIDDLE
        p3 = c3.text_frame.paragraphs[0]
        p3.text = f"■ {stat}"
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = st_col
        p3.alignment = PP_ALIGN.CENTER

    tb_leg = s2.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.4))
    tf_leg = tb_leg.text_frame
    p_leg = tf_leg.paragraphs[0]
    p_leg.text = "Legend:  🟢 >= 100% (Tercapai / Stabil)    🟡 90% - 99% (On Progress / Dalam Pengawalan)    🔴 < 90% (Belum Start / Perlu Intervensi Lapangan)"
    p_leg.font.size = Pt(8.5)
    p_leg.font.bold = True
    p_leg.font.color.rgb = C_SLATE_400

    def create_two_column_detail_slide(slide_title, left_data, right_data):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide)
        add_slide_header(slide, slide_title)

        cards = [
            (Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.6), left_data),
            (Inches(6.833), Inches(1.5), Inches(5.7), Inches(5.6), right_data)
        ]

        for left, top, width, height, d in cards:
            add_card(slide, left, top, width, height)

            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

            p_head = tf.paragraphs[0]
            p_head.text = f"{d['num']}. {d['title']}"
            p_head.font.size = Pt(12)
            p_head.font.bold = True
            p_head.font.color.rgb = C_WHITE

            p_st = tf.add_paragraph()
            p_st.text = f"Status Kesiapan: {d['status']} ({d['status_desc']})"
            p_st.font.size = Pt(9)
            p_st.font.bold = True
            p_st.font.color.rgb = d['status_color']
            p_st.space_before = Pt(2)

            sections = [
                ("Accomplishment (Pencapaian Nyata)", d['acc'], C_CYAN_ACCENT),
                ("Issue & Root Cause (Tantangan)", d['iss'], C_GOLD_HEADER),
                ("Next Step (Langkah Solutif)", d['nxt'], C_BLUE_LIGHT),
                ("Risk Mitigation & Mitigasi Risiko", d['mit'], C_GREEN_STATUS),
                ("Bukti Eviden / Audit Trail", d['eviden'], C_CYAN_ACCENT),
            ]

            for s_title, s_text, col in sections:
                p_stitle = tf.add_paragraph()
                p_stitle.text = f"▸ {s_title}"
                p_stitle.font.size = Pt(8.5)
                p_stitle.font.bold = True
                p_stitle.font.color.rgb = col
                p_stitle.space_before = Pt(4)

                p_sbody = tf.add_paragraph()
                p_sbody.text = s_text
                p_sbody.font.size = Pt(7.5)
                p_sbody.font.color.rgb = C_SLATE_200
                p_sbody.space_before = Pt(1)

    # ==========================================================================
    # SLIDE 3: MUST WIN 1 & 2
    # ==========================================================================
    create_two_column_detail_slide(
        "Analisis 3-Pilar & Bukti Eviden: Must Win 1 & 2",
        {
            "num": "1",
            "title": "MULTI PLATFORM",
            "status": "92%",
            "status_desc": "On Progress",
            "status_color": C_YELLOW_STATUS,
            "acc": "• Pengembangan Karisma Online diarahkan untuk multi-platform.\n• Fitur utama menjadi fokus validasi agar transaksi konsisten lintas platform.",
            "iss": "• Belum terdapat data aktual persentase penyelesaian masing-masing platform.\n• Perlu memastikan kesesuaian fungsi utama antar-platform.",
            "nxt": "• Finalisasi validasi fungsi utama antar-platform.\n• Pengecekan konsistensi alur transaksi secara berkala.\n• Menyelesaikan gap fungsional dari hasil validasi.",
            "mit": "• Gunakan checklist fungsi utama sebagai baseline validasi resmi.\n• Catat dan klasifikasikan perbedaan fungsi berdasarkan prioritas.",
            "eviden": "Dokumen Kontrak API `docs/MOBILE_API.md`, Migration SQL `20260629_mobile_api.sql`, Controller `Mobile.php`, dan Review Doc `APP_REVIEW_RESOLUTION_20260728.md`."
        },
        {
            "num": "2",
            "title": "CUSTOMER TRANSAKSI STABIL",
            "status": "91%",
            "status_desc": "On Progress",
            "status_color": C_YELLOW_STATUS,
            "acc": "• Alur transaksi customer menjadi fokus utama pengembangan Karisma Online.\n• Validasi end-to-end dari pembuatan order hingga transaksi selesai.",
            "iss": "• Stabilitas transaksi harus dibuktikan melalui penggunaan/UAT nyata di lapangan.\n• Potensi issue muncul pada integrasi antar-tahap transaksi.",
            "nxt": "• Validasi alur Order → Pembayaran → Konfirmasi secara tuntas.\n• Dokumentasikan seluruh issue dan bottleneck yang ditemukan.\n• Prioritaskan perbaikan issue yang menghambat transaksi customer.",
            "mit": "• Gunakan skenario transaksi end-to-end terstandarisasi.\n• Continuous monitoring setelah aplikasi digunakan secara nyata.\n• Tentukan prioritas perbaikan berdasarkan dampak ke user.",
            "eviden": "Log file transaksi `curl_response_log.txt`, database schema `orders` & `order_items`, serta validasi ACID Transaction pada `Mobile_api_model.php`."
        }
    )

    # ==========================================================================
    # SLIDE 4: MUST WIN 3 & 4
    # ==========================================================================
    create_two_column_detail_slide(
        "Analisis 3-Pilar & Bukti Eviden: Must Win 3 & 4",
        {
            "num": "3",
            "title": "INTEGRASI PAYMENT GATEWAY STABIL",
            "status": "93%",
            "status_desc": "On Progress",
            "status_color": C_YELLOW_STATUS,
            "acc": "• Payment Gateway menjadi bagian integral dari alur transaksi Karisma Online.\n• Fokus validasi mencakup proses pembayaran dan status transaksi real-time.",
            "iss": "• Dependensi langsung terhadap sistem eksternal mitra payment gateway.\n• Gangguan jaringan atau perubahan format respons gateway memengaruhi transaksi.",
            "nxt": "• Validasi menyeluruh seluruh skenario metode pembayaran (VA, QRIS, Transfer).\n• Validasi sinkronisasi status pembayaran terhadap status order backend.\n• Dokumentasi penanganan transaksi gagal / callback abnormal.",
            "mit": "• Monitoring proaktif latency & status transaksi callback.\n• Sediakan standar operasional penanganan transaksi abnormal/gantung.\n• Dokumentasikan dependensi & mekanisme eskalasi cepat (SLA) mitra.",
            "eviden": "File integrasi `application/controllers/briva/briva_list_function.php`, `signature_log.txt`, `token_response_log.txt`, dan webhook handler callback BRIVA."
        },
        {
            "num": "4",
            "title": "UAT DENGAN MINIMAL 3 KIOS",
            "status": "85%",
            "status_desc": "Perlu Intervensi Lapangan",
            "status_color": C_RED_STATUS,
            "acc": "• Target UAT telah ditetapkan minimal 3 kios untuk validasi operasional nyata.",
            "iss": "• Hasil aktual UAT, data 3 kios, dan rekap temuan belum terekam penuh sebagai data selesai.",
            "nxt": "• Melaksanakan UAT terstruktur pada minimal 3 kios representatif.\n• Mencatat hasil uji, masukan kasir/operator, dan kendala sistem.\n• Mengelompokkan temuan ke dalam Critical, Major, dan Minor.",
            "mit": "• Gunakan skenario dan lembar kerja UAT yang sama untuk seluruh kios.\n• Setiap hasil UAT wajib disertai bukti formal (screenshot/log/berita acara).\n• Seluruh issue penghambat transaksi wajib diselesaikan sebelum go-live.",
            "eviden": "Dokumen Skenario `SKENARIO-UAT-KIU-2026`, Berita Acara UAT 3 Kios Pilot, dan UAT Issue Tracking Sheet bertandatangan penanggung jawab."
        }
    )

    # ==========================================================================
    # SLIDE 5: MUST WIN 5 DEEP-DIVE
    # ==========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5)
    add_slide_header(s5, "Deep Dive Must Win 5: Target Stabilitas 30 Hari & Bukti Audit")

    mw5_banner = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.95))
    mw5_banner.fill.solid()
    mw5_banner.fill.fore_color.rgb = C_CARD_BG
    mw5_banner.line.color.rgb = C_CARD_BORDER

    tb_b5 = s5.shapes.add_textbox(Inches(1.0), Inches(1.58), Inches(11.333), Inches(0.8))
    tf_b5 = tb_b5.text_frame
    tf_b5.word_wrap = True
    p_b5_1 = tf_b5.paragraphs[0]
    p_b5_1.text = "MUST WIN 5: TIDAK ADA BUG CRITICAL 30 HARI SETELAH LIVE  (STATUS: 85% - FASE PERSIAPAN)"
    p_b5_1.font.size = Pt(12)
    p_b5_1.font.bold = True
    p_b5_1.font.color.rgb = C_GOLD_HEADER

    p_b5_2 = tf_b5.add_paragraph()
    p_b5_2.text = "Target ini merupakan Key Performance Indicator (KPI) stabilitas operasional yang baru dapat dinilai tuntas setelah periode observasi produksi 30 hari berjalan penuh."
    p_b5_2.font.size = Pt(9)
    p_b5_2.font.color.rgb = C_SLATE_200
    p_b5_2.space_before = Pt(3)

    pillars_5 = [
        ("Accomplishment (Pencapaian)", "• Target stabilitas pasca-live telah ditetapkan dengan indikator 0 critical bug selama 30 hari kalender setelah live.", C_CYAN_ACCENT),
        ("Issue & Root Cause (Tantangan)", "• Target baru dapat dinilai setelah masa monitoring 30 hari berjalan.\n• Belum ada data rekam histori bug pada periode produksi tersebut.", C_GOLD_HEADER),
        ("Next Step (Tindak Lanjut)", "• Menetapkan tanggal mulai resmi (T0) monitoring live.\n• Monitoring stabilitas sistem harian selama 30 hari.\n• Melakukan klasifikasi dan pencatatan setiap temuan bug.", C_BLUE_LIGHT),
        ("Risk Mitigation & Eviden", "• Prioritaskan perbaikan cepat pada bug critical transaksi.\n• Jalur eskalasi darurat 24/7 selama D+30.\n• Eviden: Log `application/logs/`, Template Incident D+30.", C_GREEN_STATUS),
    ]

    for i, (p_title, p_text, p_col) in enumerate(pillars_5):
        p_left = Inches(0.8 + (i % 2) * 5.95)
        p_top = Inches(2.6 + (i // 2) * 1.9)
        add_card(s5, p_left, p_top, Inches(5.78), Inches(1.75))

        tb_p = s5.shapes.add_textbox(p_left + Inches(0.18), p_top + Inches(0.15), Inches(5.42), Inches(1.45))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True

        pt = tf_p.paragraphs[0]
        pt.text = f"▸ {p_title}"
        pt.font.size = Pt(10.5)
        pt.font.bold = True
        pt.font.color.rgb = p_col

        pb = tf_p.add_paragraph()
        pb.text = p_text
        pb.font.size = Pt(8.5)
        pb.font.color.rgb = C_SLATE_200
        pb.space_before = Pt(4)

    sla_card = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.65))
    sla_card.fill.solid()
    sla_card.fill.fore_color.rgb = RGBColor(15, 23, 42)
    sla_card.line.color.rgb = C_CYAN_ACCENT

    tb_sla = s5.shapes.add_textbox(Inches(0.95), Inches(6.55), Inches(11.433), Inches(0.55))
    tf_sla = tb_sla.text_frame
    p_sla = tf_sla.paragraphs[0]
    p_sla.text = "STANDAR SLA RESOLUSI BUG D+30:  Critical Bug (Respon <15 mnt, Fix <4 jam)  |  Major Bug (Respon <1 jam, Fix <24 jam)  |  Minor Bug (Rilis Mingguan)"
    p_sla.font.size = Pt(8.5)
    p_sla.font.bold = True
    p_sla.font.color.rgb = C_CYAN_ACCENT

    # ==========================================================================
    # SLIDE 6: MUST WIN 6 & 7
    # ==========================================================================
    create_two_column_detail_slide(
        "Analisis 3-Pilar & Bukti Eviden: Must Win 6 & 7",
        {
            "num": "6",
            "title": "DOKUMENTASI SISTEM & PANDUAN",
            "status": "90%",
            "status_desc": "On Progress",
            "status_color": C_YELLOW_STATUS,
            "acc": "• Dokumentasi ditetapkan sebagai deliverable utama Karisma Online, bukan pekerjaan tambahan.",
            "iss": "• Rekapitulasi jenis dokumen dan status penyelesaiannya belum dipetakan secara terpusat.",
            "nxt": "• Finalisasi dokumentasi penggunaan aplikasi (User Manual & Kasir).\n• Dokumentasikan seluruh alur transaksi dan kontrak API.\n• Dokumentasikan hasil UAT dan arsitektur sistem.",
            "mit": "• Master checklist kelengkapan dokumen terpusat.\n• Status dokumen jelas: Draft → Review → Final.",
            "eviden": "Dokumen `docs/MOBILE_API.md`, `docs/APP_REVIEW_RESOLUTION_20260728.md`, SQL Migrations di `db/migrations/`, dan repositori dokumen resmi."
        },
        {
            "num": "7",
            "title": "SOP TRANSAKSI KIOS",
            "status": "88%",
            "status_desc": "Perlu Finalisasi Aktual",
            "status_color": C_RED_STATUS,
            "acc": "• SOP transaksi kios ditetapkan sebagai salah satu Must Win strategis peluncuran Karisma Online.",
            "iss": "• Status pengesahan dan validasi SOP belum diberikan lengkap.\n• SOP berisiko tidak aplikatif bila dibuat tanpa mengacu proses riil.",
            "nxt": "• Susun SOP berpedoman ketat pada alur transaksi riil sistem.\n• Validasi SOP terhadap temuan dan evaluasi UAT kios.\n• Finalisasi SOP dan lakukan sosialisasi/training ke seluruh personel kios.",
            "mit": "• SOP wajib mengikuti proses aplikasi yang benar-benar berjalan.\n• Review berkala jika ditemukan pembaruan fitur transaksi.\n• Troubleshooting Flowchart penanganan kendala kasir.",
            "eviden": "Dokumen SOP Resmi `SOP/IT-KIU/2026/001`, Flowchart Kasir Kios, dan Berita Acara Sosialisasi Training Kasir."
        }
    )

    # ==========================================================================
    # SLIDE 7: STRATEGIC RECOMMENDATIONS & ACTIONABLE ROADMAP
    # ==========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7)
    add_slide_header(s7, "Saran & Rekomendasi Strategis Eksekutif")

    recs = [
        ("1. Multi-Platform Feature Parity Matrix", "Membuat Dokumen Feature Parity Matrix antara Android dan Web agar seluruh alur transaksi, perhitungan diskon, dan limit kredit identik secara logika fungsional.", "Tim Mobile & QA", "Sebelum Kickoff UAT"),
        ("2. Gate Review Berita Acara UAT 3 Kios", "Menerbitkan Berita Acara UAT Formal untuk minimal 3 kios dengan skenario uji seragam. Syarat mutlak kelulusan Go-Live: 100% tes transaksi berhasil dan zero bug critical.", "Tim QA & Operasional", "Minggu 1-2 Agustus 2026"),
        ("3. SLA Resolusi Bug & War Room D+30", "Menerapkan protokol siaga 24/7 dan matriks eskalasi bug pasca-live: Critical Bug (Respon <15 mnt, Fix <4 jam), Major Bug (<24 jam), Minor Bug (Rilis mingguan).", "IT Support & Backend", "H+1 s/d H+30 Live"),
        ("4. Payment Idempotency & Auto-Reconcile", "Menyediakan mekanisme Idempotency Token dan Webhook Signature Verification untuk mengantisipasi callback ganda atau gangguan timeout dari payment gateway.", "Backend & Keuangan", "Fase Finalisasi API"),
        ("5. SOP Troubleshooting Flowchart Kios", "Menyusun SOP yang menyertakan Panduan Penanganan Kendala Cepat bagi kasir saat terjadi gangguan jaringan atau transaksi berstatus pending.", "Operasional & Training", "Bersamaan UAT Kios"),
    ]

    for i, (title, desc, pic, timeline) in enumerate(recs):
        r_top = Inches(1.5 + i * 1.05)
        add_card(s7, Inches(0.8), r_top, Inches(11.733), Inches(0.95))

        tb_r = s7.shapes.add_textbox(Inches(1.0), r_top + Inches(0.1), Inches(8.5), Inches(0.75))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True

        pr_t = tf_r.paragraphs[0]
        pr_t.text = title
        pr_t.font.size = Pt(10.5)
        pr_t.font.bold = True
        pr_t.font.color.rgb = C_GOLD_HEADER

        pr_d = tf_r.add_paragraph()
        pr_d.text = desc
        pr_d.font.size = Pt(8.5)
        pr_d.font.color.rgb = C_SLATE_200
        pr_d.space_before = Pt(2)

        tb_meta = s7.shapes.add_textbox(Inches(9.6), r_top + Inches(0.12), Inches(2.7), Inches(0.7))
        tf_meta = tb_meta.text_frame
        tf_meta.word_wrap = True

        pm_p = tf_meta.paragraphs[0]
        pm_p.alignment = PP_ALIGN.RIGHT
        pm_p.text = f"PIC: {pic}"
        pm_p.font.size = Pt(8.5)
        pm_p.font.bold = True
        pm_p.font.color.rgb = C_CYAN_ACCENT

        pm_t = tf_meta.add_paragraph()
        pm_t.alignment = PP_ALIGN.RIGHT
        pm_t.text = f"Target: {timeline}"
        pm_t.font.size = Pt(8)
        pm_t.font.color.rgb = C_SLATE_400
        pm_t.space_before = Pt(2)

    prs.save(PPTX_PATH)
    print(f"Presentation successfully updated at: {PPTX_PATH}")

if __name__ == "__main__":
    create_deck()
