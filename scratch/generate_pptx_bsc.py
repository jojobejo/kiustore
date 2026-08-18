import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]

# Color Palette (Obsidian Deep Navy Futuristic Executive Theme)
BG_COLOR = RGBColor(11, 19, 43)        # #0B132B
CARD_BG = RGBColor(28, 37, 65)         # #1C2541
CARD_BORDER = RGBColor(58, 80, 107)    # #3A506B
CYAN = RGBColor(0, 240, 255)           # #00F0FF
CYAN_DARK = RGBColor(14, 165, 233)     # #0EA5E9
WHITE = RGBColor(255, 255, 255)
GRAY_TEXT = RGBColor(148, 163, 184)    # #94A3B8
EMERALD = RGBColor(16, 185, 129)       # #10B981
AMBER = RGBColor(245, 158, 11)         # #F59E0B
ROSE = RGBColor(239, 68, 68)           # #EF4444
DARK_BG = RGBColor(15, 23, 42)

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text, category_text="BALANCED SCORECARD & EXECUTIVE REPORT"):
    # Header container
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = CYAN
    p0.font.name = "Calibri"
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Calibri"

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

# ==============================================================================
# SLIDE 1: COVER & HERO TITLE
# ==============================================================================
s1 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s1)

# Hero Card
add_card(s1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), bg_color=CARD_BG, border_color=CYAN)

tb1 = s1.shapes.add_textbox(Inches(1.3), Inches(1.3), Inches(10.733), Inches(4.8))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "PT. KARISMA INDOAGRO UNIVERSAL"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CYAN

p = tf1.add_paragraph()
p.text = "KARISMA ONLINE"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

p = tf1.add_paragraph()
p.text = "Balanced Scorecard (BSC) & Strategic Multi-Platform Launching Report"
p.font.size = Pt(18)
p.font.color.rgb = GRAY_TEXT

p = tf1.add_paragraph()
p.text = "\nEksplorasi Kesiapan Ekosistem Digital Multi-Platform (iOS Swift Native, Android Kotlin/Java, CI3 RESTful API PHP 8.x) Terintegrasi Fintech BRIVA SNAP & Logistik Presisi."
p.font.size = Pt(13)
p.font.color.rgb = WHITE

p = tf1.add_paragraph()
p.text = "\n🟢 STATUS KELAYAKAN SISTEM: 96.79% (LAUNCH READY)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = EMERALD

p = tf1.add_paragraph()
p.text = "Periode Evaluasi: Q3 2026 | Dokumen Resmi Eksekutif & Dewan Direksi"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = GRAY_TEXT

# ==============================================================================
# SLIDE 2: EXECUTIVE SUMMARY & BSC PROJECT FLOW MAPPING
# ==============================================================================
s2 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s2)
add_header(s2, "Executive Summary & Pemetaan BSC Sesuai Alur Proyek", "IKHTISAR STRATEGIS")

cards_s2 = [
    ("1. ALUR FONDASI & SISTEM", "Learning & Growth (Bobot: 20%)", "🟢 98.75%", "Arsitektur Tri-Platform Native (iOS/Android/CI3), Token SHA-256, BCRYPT, Isolasi Akun Demo.", EMERALD),
    ("2. ALUR RANTAI PASOK", "Internal Process (Bobot: 30%)", "🟡 92.40%", "Logistik Presisi RajaOngkir Pro (Kecamatan), 30-Min Quote Lock, Dual-Unit Stok & Deadstock Engine.", AMBER),
    ("3. ALUR PASAR & KIOS", "Customer & Market (Bobot: 25%)", "🟢 100.0%", "Apple Guideline 5.1.1(v) In-App Deletion, Guest Browsing, CS Live Chat, Android Play Store V1.", EMERALD),
    ("4. ALUR FINANSIAL & KAS", "Financial Perspective (Bobot: 25%)", "🟢 98.00%", "BRI BRIVA SNAP 15-Min VA, Server-Side Multi-Tier Margin Guard, Fast Cash Flow DSO Acceleration.", EMERALD),
]

for idx, (title, sub, score, desc, color) in enumerate(cards_s2):
    left = Inches(0.8 + idx * 2.98)
    add_card(s2, left, Inches(1.5), Inches(2.8), Inches(4.3))
    
    tb = s2.shapes.add_textbox(left + Inches(0.15), Inches(1.65), Inches(2.5), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    p = tf.add_paragraph()
    p.text = sub
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT
    
    p = tf.add_paragraph()
    p.text = f"\n{score}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    
    p = tf.add_paragraph()
    p.text = f"\n{desc}"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE

# Total bar
add_card(s2, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9), bg_color=CARD_BG, border_color=EMERALD)
tb_tot = s2.shapes.add_textbox(Inches(1.0), Inches(6.05), Inches(11.333), Inches(0.8))
tf_tot = tb_tot.text_frame
p = tf_tot.paragraphs[0]
p.text = "SKOR KESIAPAN KESELURUHAN (WEIGHTED COMPOSITE SCORE): 96.79% 🟢 LAUNCH READY"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = EMERALD
p = tf_tot.add_paragraph()
p.text = "Seluruh modul inti lulus audit regulasi & integrasi fintech. Proyek siap dirilis serentak ke publik."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 3: BSC PERSPECTIVE 1 - LEARNING & GROWTH
# ==============================================================================
s3 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s3)
add_header(s3, "Perspektif 1: Learning, Growth & System Architecture", "ALUR FONDASI SISTEM (BOBOT: 20% | SKOR: 98.75% 🟢)")

# Left Card: Table KPIs
add_card(s3, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.4))
tb3_l = s3.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(7.1), Inches(5.1))
tf3_l = tb3_l.text_frame
tf3_l.word_wrap = True

p = tf3_l.paragraphs[0]
p.text = "MATRIKS SASARAN STRATEGIS (MUST-WIN KPIs)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN

kpis_s3 = [
    ("LRN-01: Tri-Platform Native Architecture", "Swift iOS + Kotlin Android + CI3 HMVC PHP 8.x", "🟢 100%"),
    ("LRN-02: Keamanan & Autentikasi Modern", "Bearer Token SHA-256 & Hashing BCRYPT (mobile_api_tokens)", "🟢 100%"),
    ("LRN-03: Manajemen Wilayah Salesman", "Pemetaan Kios Binaan per Salesman Lapangan", "🟢 95%"),
    ("LRN-04: Isolasi Akun Demo Internal", "Flagging users.is_internal Anti-Distorsi Omset", "🟢 100%"),
]
for code, desc, stat in kpis_s3:
    p = tf3_l.add_paragraph()
    p.text = f"\n• {code} [{stat}]"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf3_l.add_paragraph()
    p.text = f"   {desc}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT

# Right Card: 3-Pilar Grounded Analysis
add_card(s3, Inches(8.5), Inches(1.5), Inches(4.033), Inches(5.4), border_color=CYAN_DARK)
tb3_r = s3.shapes.add_textbox(Inches(8.7), Inches(1.65), Inches(3.633), Inches(5.1))
tf3_r = tb3_r.text_frame
tf3_r.word_wrap = True

p = tf3_r.paragraphs[0]
p.text = "3-PILAR ANALISIS STRATEGIS"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN

p = tf3_r.add_paragraph()
p.text = "\n1. Accomplishment (Pencapaian):"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = EMERALD
p = tf3_r.add_paragraph()
p.text = "Backend CI3 HMVC terisolasi penuh pada /api/v1, throughput tinggi dan memory-efficient."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

p = tf3_r.add_paragraph()
p.text = "\n2. Issues & Root Cause (Tantangan):"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = AMBER
p = tf3_r.add_paragraph()
p.text = "Keterbatasan sinyal internet 3G di pelosok persawahan saat salesman sinkronisasi data."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

p = tf3_r.add_paragraph()
p.text = "\n3. Next Steps & Mitigasi:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = CYAN
p = tf3_r.add_paragraph()
p.text = "Optimasi payload JSON < 50KB dan local SQLite caching untuk mode semi-offline."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 4: BSC PERSPECTIVE 2 - INTERNAL PROCESS
# ==============================================================================
s4 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s4)
add_header(s4, "Perspektif 2: Internal Business Process & Logistics", "ALUR RANTAI PASOK & OPERASIONAL (BOBOT: 30% | SKOR: 92.40% 🟡)")

# Left Card
add_card(s4, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.4))
tb4_l = s4.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(7.1), Inches(5.1))
tf4_l = tb4_l.text_frame
tf4_l.word_wrap = True

p = tf4_l.paragraphs[0]
p.text = "MATRIKS SASARAN STRATEGIS (MUST-WIN KPIs)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN

kpis_s4 = [
    ("PRC-01: Logistik Presisi RajaOngkir Pro", "Kalkulasi tarif kurir hingga Subdistrict ID (Kecamatan)", "🟢 100%"),
    ("PRC-02: 30-Min Shipping Quote Lock", "Penguncian kuotasi tarif pada tabel mobile_shipping_quotes", "🟢 100%"),
    ("PRC-03: Dual-Unit Gramasi & Stok Real-Time", "Otomasi konversi Botol/Pcs vs Dus/Karton multiplier", "🟢 100%"),
    ("PRC-04: Deadstock Analytics Engine", "Deteksi produk slow-moving > 1 tahun (penjualan < 5%)", "🟡 92%"),
    ("PRC-05: Transaksi Kredit & Limit Tempo", "Integrasi limit kredit web ke antarmuka mobile", "🔴 78%"),
]
for code, desc, stat in kpis_s4:
    p = tf4_l.add_paragraph()
    p.text = f"\n• {code} [{stat}]"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf4_l.add_paragraph()
    p.text = f"   {desc}"
    p.font.size = Pt(9.5)
    p.font.color.rgb = GRAY_TEXT

# Right Card
add_card(s4, Inches(8.5), Inches(1.5), Inches(4.033), Inches(5.4), border_color=AMBER)
tb4_r = s4.shapes.add_textbox(Inches(8.7), Inches(1.65), Inches(3.633), Inches(5.1))
tf4_r = tb4_r.text_frame
tf4_r.word_wrap = True

p = tf4_r.paragraphs[0]
p.text = "3-PILAR ANALISIS STRATEGIS"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN

p = tf4_r.add_paragraph()
p.text = "\n1. Accomplishment (Pencapaian):"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = EMERALD
p = tf4_r.add_paragraph()
p.text = "Menghilangkan risiko selisih ongkir berkat tarif kecamatan presisi & penguncian 30 menit."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

p = tf4_r.add_paragraph()
p.text = "\n2. Issues & Root Cause (Tantangan):"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = AMBER
p = tf4_r.add_paragraph()
p.text = "Plafon kredit tempo mobile tertahan karena verifikasi jaminan fisik masih butuh approval komite."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

p = tf4_r.add_paragraph()
p.text = "\n3. Next Steps & Mitigasi:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = CYAN
p = tf4_r.add_paragraph()
p.text = "Rilis credit scoring otomatis berbasis riwayat transaksi pada Q1 2027."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 5: BSC PERSPECTIVE 3 - CUSTOMER & MARKET
# ==============================================================================
s5 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s5)
add_header(s5, "Perspektif 3: Customer, Market Access & Compliance", "ALUR PASAR & KIOS MITRA (BOBOT: 25% | SKOR: 100.0% 🟢)")

# Left Card
add_card(s5, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.4))
tb5_l = s5.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(7.1), Inches(5.1))
tf5_l = tb5_l.text_frame
tf5_l.word_wrap = True

p = tf5_l.paragraphs[0]
p.text = "MATRIKS SASARAN STRATEGIS (MUST-WIN KPIs)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN

kpis_s5 = [
    ("CUS-01: Apple Guideline 5.1.1(v) Compliance", "Endpoint DELETE /api/v1/account & Non-PII Hash", "🟢 100%"),
    ("CUS-02: Guest Browsing Experience", "Katalog produk, kategori & promo terbuka tanpa wajib login", "🟢 100%"),
    ("CUS-03: Layanan Bantuan Terpadu (Live Chat)", "Komunikasi customer-ke-admin real-time via API messages", "🟢 100%"),
    ("CUS-04: Kesiapan Google Play Store", "Android Native Kotlin/Java MVVM & Data-Light Architecture", "🟢 100%"),
]
for code, desc, stat in kpis_s5:
    p = tf5_l.add_paragraph()
    p.text = f"\n• {code} [{stat}]"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf5_l.add_paragraph()
    p.text = f"   {desc}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT

# Right Card
add_card(s5, Inches(8.5), Inches(1.5), Inches(4.033), Inches(5.4), border_color=EMERALD)
tb5_r = s5.shapes.add_textbox(Inches(8.7), Inches(1.65), Inches(3.633), Inches(5.1))
tf5_r = tb5_r.text_frame
tf5_r.word_wrap = True

p = tf5_r.paragraphs[0]
p.text = "3-PILAR ANALISIS STRATEGIS"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN

p = tf5_r.add_paragraph()
p.text = "\n1. Accomplishment (Pencapaian):"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = EMERALD
p = tf5_r.add_paragraph()
p.text = "Lulus 100% persyaratan Apple App Store Review & Google Play Developer Policy."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

p = tf5_r.add_paragraph()
p.text = "\n2. Issues & Root Cause (Tantangan):"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = AMBER
p = tf5_r.add_paragraph()
p.text = "Resistensi kebiasaan kios konvensional yang terbiasa pesan manual via telepon/WhatsApp."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

p = tf5_r.add_paragraph()
p.text = "\n3. Next Steps & Mitigasi:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = CYAN
p = tf5_r.add_paragraph()
p.text = "Program insentif onboarding dan pendampingan salesman dengan booklet interaktif."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 6: BSC PERSPECTIVE 4 - FINANCIAL PERSPECTIVE
# ==============================================================================
s6 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s6)
add_header(s6, "Perspektif 4: Financial Revenue & Cash Flow Settlement", "ALUR FINANSIAL & MONETISASI (BOBOT: 25% | SKOR: 98.00% 🟢)")

# Left Card
add_card(s6, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.4))
tb6_l = s6.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(7.1), Inches(5.1))
tf6_l = tb6_l.text_frame
tf6_l.word_wrap = True

p = tf6_l.paragraphs[0]
p.text = "MATRIKS SASARAN STRATEGIS (MUST-WIN KPIs)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN

kpis_s6 = [
    ("FIN-01: Integrasi Bank BRI BRIVA SNAP VA", "Dynamic VA 15-Minute Expiry & Instant Settlement (briva_api)", "🟢 100%"),
    ("FIN-02: Proteksi Multi-Tier Pricing Margin", "Enkripsi Server-Side 3 Level Harga (Retail, Grosir, Distr)", "🟢 100%"),
    ("FIN-03: Otomasi Multi-Bank Manual Transfer", "Validasi MIME & Secure Receipt Hash Storage (BCA, Mandiri, BRI)", "🟢 100%"),
    ("FIN-04: Akselerasi Cash Flow & DSO", "Percepatan penerimaan kas & reduksi DSO vs siklus manual", "🟡 92%"),
]
for code, desc, stat in kpis_s6:
    p = tf6_l.add_paragraph()
    p.text = f"\n• {code} [{stat}]"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf6_l.add_paragraph()
    p.text = f"   {desc}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT

# Right Card
add_card(s6, Inches(8.5), Inches(1.5), Inches(4.033), Inches(5.4), border_color=EMERALD)
tb6_r = s6.shapes.add_textbox(Inches(8.7), Inches(1.65), Inches(3.633), Inches(5.1))
tf6_r = tb6_r.text_frame
tf6_r.word_wrap = True

p = tf6_r.paragraphs[0]
p.text = "3-PILAR ANALISIS STRATEGIS"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN

p = tf6_r.add_paragraph()
p.text = "\n1. Accomplishment (Pencapaian):"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = EMERALD
p = tf6_r.add_paragraph()
p.text = "Otomasi BRIVA memangkas siklus rekonsiliasi pembayaran menjadi instan."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

p = tf6_r.add_paragraph()
p.text = "\n2. Issues & Root Cause (Tantangan):"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = AMBER
p = tf6_r.add_paragraph()
p.text = "Kios pelosok terkendala limit harian kartu debit saat pembayaran order nominal besar."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

p = tf6_r.add_paragraph()
p.text = "\n3. Next Steps & Mitigasi:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = CYAN
p = tf6_r.add_paragraph()
p.text = "Sediakan fallback transfer multi-bank dengan antrean verifikasi kilat tim finance."
p.font.size = Pt(10)
p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 7: GROUNDED PROBLEM-SOLVING MATRIX
# ==============================================================================
s7 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s7)
add_header(s7, "3-Pilar Matriks Penyelesaian Masalah & Mitigasi Risiko", "PROBLEM-SOLVING & RISK MITIGATION MATRIX")

rows_s7 = [
    ("1. Regulasi Apple App Store", "Lulus Apple Guideline 5.1.1(v) via Account Deletion API & Guest Browsing.", "Data PII di-hash SHA-256, faktur dianonimkan untuk pajak.", "Tim iOS & Backend", "🟢 Done"),
    ("2. Fintech & Pembayaran VA", "Integrasi BRI BRIVA SNAP 15-min auto-settlement real-time.", "Fallback transfer manual multi-bank dengan fast-track verifikasi admin.", "Tim Finance & IT", "🟢 Done"),
    ("3. Logistik Pedesaan & Armada", "RajaOngkir Pro tingkat kecamatan & 30-min price locking.", "Integrasikan opsi armada truk internal untuk tonase besar di modul admin.", "Tim Logistik & Gudang", "🟡 Q4 2026"),
    ("4. Deadstock Management", "Deteksi produk slow-moving > 1 tahun dengan penjualan < 5%.", "Hubungkan data deadstock ke modul flash sale clearance otomatis.", "Supply Chain & Sales", "🟡 Q4 2026"),
    ("5. Limit Kredit Tempo Mobile", "Modul Piutang & limit kredit berjalan stabil di Web Enterprise.", "Kembangkan scoring kredit otomatis berbasis riwayat sebelum buka di mobile.", "Finance & Credit Comm", "🔴 Q1 2027"),
]

for idx, (dom, acc, mit, pic, stat) in enumerate(rows_s7):
    top = Inches(1.5 + idx * 1.1)
    add_card(s7, Inches(0.8), top, Inches(11.733), Inches(0.95))
    tb = s7.shapes.add_textbox(Inches(1.0), top + Inches(0.08), Inches(11.333), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"{dom}  |  PIC: {pic}  |  Status: {stat}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN if "Done" not in stat else EMERALD
    
    p = tf.add_paragraph()
    p.text = f"Capaian: {acc}  →  Mitigasi: {mit}"
    p.font.size = Pt(9.5)
    p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 8: TECH ARCHITECTURE & SECURITY
# ==============================================================================
s8 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s8)
add_header(s8, "Arsitektur Teknologi Tri-Platform & Keamanan Terintegrasi", "INFRASTRUKTUR & KEAMANAN SISTEM")

tech_cards = [
    ("iOS NATIVE APP", "Swift & SwiftUI\n\n• HIG Compliant 100%\n• Secure Keychain Auth\n• iPad Adaptive UI\n• In-App Deletion Flow\n• Guest Browsing Mode", CYAN),
    ("ANDROID NATIVE APP", "Kotlin / Java MVVM\n\n• Material Design 3\n• Lightweight APK Build\n• Data-Light Architecture\n• Offline Image Cache\n• Live Support Chat", EMERALD),
    ("ENTERPRISE BACKEND", "CodeIgniter 3 PHP 8.x\n\n• HMVC Modular Architecture\n• RESTful JSON API v1\n• Bearer SHA-256 Token\n• BCRYPT Password Hash\n• ACID Transaction Locks", CYAN_DARK),
    ("DATABASE & INTEGRASI", "MariaDB & External APIs\n\n• InnoDB ACID Transactions\n• Bank BRI BRIVA SNAP VA\n• RajaOngkir Pro API\n• Non-PII Audit Deletion\n• Multi-Tier Price Views", AMBER),
]

for idx, (title, content, color) in enumerate(tech_cards):
    left = Inches(0.8 + idx * 2.98)
    add_card(s8, left, Inches(1.5), Inches(2.8), Inches(5.4), border_color=color)
    tb = s8.shapes.add_textbox(left + Inches(0.15), Inches(1.7), Inches(2.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    
    p = tf.add_paragraph()
    p.text = f"\n{content}"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 9: WEB ENTERPRISE BACK-OFFICE
# ==============================================================================
s9 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s9)
add_header(s9, "Web Enterprise Back-Office & Manajemen Operasional", "TATA KELOLA OPERASIONAL PUSAT")

admin_cards = [
    ("Dashboard Eksekutif", "Visualisasi tren omset, rasio konversi pesanan, dan monitoring piutang salesman secara real-time.", CYAN),
    ("Antrean Packing & Pengiriman", "Workflow pemrosesan order gudang, cetak surat jalan, faktur pajak, dan integrasi nomor resi kurir.", EMERALD),
    ("Sales Force & Territory", "Pemetaan kios binaan per salesman lapangan untuk monitoring target penjualan dan rute kunjungan.", CYAN_DARK),
    ("Verifikasi Keuangan & Piutang", "Rekonsiliasi otomatis BRIVA dan validasi cepat bukti transfer manual multi-bank.", AMBER),
]

for idx, (title, desc, color) in enumerate(admin_cards):
    row = idx // 2
    col = idx % 2
    left = Inches(0.8 + col * 5.96)
    top = Inches(1.6 + row * 2.7)
    add_card(s9, left, top, Inches(5.76), Inches(2.4), border_color=color)
    tb = s9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.36), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"• {title}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    
    p = tf.add_paragraph()
    p.text = f"\n{desc}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 10: GO-TO-MARKET ROLLOUT
# ==============================================================================
s10 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s10)
add_header(s10, "Rencana Peluncuran & Strategi Go-To-Market (3 Tahap)", "STRATEGI ROLLOUT & IMPLEMENTASI")

phases = [
    ("FASE 1: STORE RELEASE", "Agustus 2026", "Publikasi serentak di Apple App Store & Google Play Store, pengujian live production end-to-end.", CYAN),
    ("FASE 2: PILOT ONBOARDING", "September 2026", "Onboarding 100 Kios Mitra percontohan di area Jawa Timur & Jawa Tengah dengan pendampingan salesman.", EMERALD),
    ("FASE 3: NATIONAL ROLLOUT", "Oktober 2026 - Seterusnya", "Rollout penuh ke seluruh jaringan kios nasional, integrasi program reward poin dan kampanye promo regional.", AMBER),
]

for idx, (phase, date, desc, color) in enumerate(phases):
    left = Inches(0.8 + idx * 3.98)
    add_card(s10, left, Inches(1.6), Inches(3.76), Inches(5.2), border_color=color)
    tb = s10.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.36), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    
    p = tf.add_paragraph()
    p.text = date
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GRAY_TEXT
    
    p = tf.add_paragraph()
    p.text = f"\n{desc}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

# ==============================================================================
# SLIDE 11: INNOVATION ROADMAP
# ==============================================================================
s11 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s11)
add_header(s11, "Roadmap Inovasi Berkelanjutan (2026 - 2027)", "ROADMAP TEKNOLOGI & INOVASI")

roadmaps = [
    ("Q4 2026", "Deadstock Clearance & Internal Fleet", "Otomasi flash sale obat slow moving dan penambahan opsi pengiriman armada truk internal.", CYAN),
    ("Q1 2027", "Smart Credit Scoring & Plafon Tempo", "Algoritma penilaian kelayakan kredit otomatis berbasis histori transaksi untuk belanja tempo mobile.", EMERALD),
    ("Q2 2027", "AI Crop Disease Diagnosis", "Rekomendasi produk fungisida/insektisida berbasis scan foto daun & deteksi hama berbasis AI.", AMBER),
]

for idx, (quarter, title, desc, color) in enumerate(roadmaps):
    left = Inches(0.8 + idx * 3.98)
    add_card(s11, left, Inches(1.6), Inches(3.76), Inches(5.2), border_color=color)
    tb = s11.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.36), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = quarter
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = f"\n{desc}"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY_TEXT

# ==============================================================================
# SLIDE 12: CONCLUSION & EXECUTIVE RECOMMENDATIONS
# ==============================================================================
s12 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s12)
add_header(s12, "Kesimpulan Eksekutif & Rekomendasi Dewan Direksi", "PENUTUP & KEPUTUSAN STRATEGIS")

add_card(s12, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4), bg_color=CARD_BG, border_color=EMERALD)
tb12 = s12.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.933), Inches(4.8))
tf12 = tb12.text_frame
tf12.word_wrap = True

p = tf12.paragraphs[0]
p.text = "RINGKASAN STATUS KESIAPAN AKHIR (FINAL READINESS VERIFICATION):"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CYAN

points12 = [
    "Ekosistem digital multi-platform Karisma Online telah mencapai skor kesiapan terverifikasi 96.79% (🟢 LAUNCH READY).",
    "Seluruh hambatan regulasi Apple App Store Guideline 5.1.1(v) telah diselesaikan 100% dengan mekanisme penghapusan akun mandiri dan hashing non-PII.",
    "Integrasi fintech BRI BRIVA SNAP telah aktif dengan sistem dynamic VA 15 menit dan rekonsiliasi instan.",
    "Logistik presisi RajaOngkir Pro telah mengunci tarif hingga tingkat kecamatan tujuan.",
    "Direkomendasikan kepada Dewan Direksi untuk menyetujui peluncuran publik (Go-Live) dan pelaksanaan Pilot Project Onboarding 100 Kios Mitra."
]
for pt in points12:
    p = tf12.add_paragraph()
    p.text = f"\n✓  {pt}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

prs.save("docs/PRESENTASI_KARISMA_ONLINE.pptx")
print("Successfully generated docs/PRESENTASI_KARISMA_ONLINE.pptx")
